"""Update the existing Level 2 Capstone PPTX matching the exact pattern, formatting, and requirements."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def update_presentation():
    pptx_source = r"C:\Users\7PIN\Desktop\ICAI_AI_Course\Level_2\Capstone\ICAI_AI_Course_Level_2_Presentation_Updated.pptx"
    prs = Presentation(pptx_source)
    
    # --------------------------------------------------------------------------
    # SLIDE 1: Title Slide
    # --------------------------------------------------------------------------
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.name == "Rounded Rectangle 3":
            shape.text_frame.paragraphs[0].text = "ICAI AI COURSE — LEVEL 2 • CAPSTONE PROJECT"
        elif shape.name == "TextBox 4":
            tf = shape.text_frame
            tf.paragraphs[0].text = "Schedule III Ratio Analyser"
            if len(tf.paragraphs) > 1:
                tf.paragraphs[1].text = "Automated Statutory 11 Mandated Financial Ratios & Variance Analysis (Companies Act, 2013)"
        elif shape.name == "Rounded Rectangle 6":
            tf = shape.text_frame
            tf.paragraphs[0].text = "CAPSTONE PROJECT"
            if len(tf.paragraphs) > 1:
                tf.paragraphs[1].text = "Schedule III Ratio Analyser"
            if len(tf.paragraphs) > 2:
                tf.paragraphs[2].text = "Statutory Audit Intelligence Engine"
        elif shape.name == "Rounded Rectangle 7":
            tf = shape.text_frame
            tf.paragraphs[0].text = "SUBMISSION DETAILS"
        elif shape.name == "TextBox 8":
            tf = shape.text_frame
            tf.paragraphs[0].text = "• Submitted By : xxx"
            if len(tf.paragraphs) > 1:
                tf.paragraphs[1].text = "• ICAI AI Course Level 2"
            if len(tf.paragraphs) > 2:
                tf.paragraphs[2].text = "• Desktop Offline Application"
        elif shape.name == "TextBox 9":
            tf = shape.text_frame
            tf.paragraphs[0].text = "• Standalone Executable (.exe)"
            if len(tf.paragraphs) > 1:
                tf.paragraphs[1].text = "• 100% Offline & Secure"

    # --------------------------------------------------------------------------
    # SLIDE 2: Problem Statement
    # --------------------------------------------------------------------------
    s2 = prs.slides[1]
    for shape in s2.shapes:
        if shape.name == "TextBox 3":
            tf = shape.text_frame
            tf.paragraphs[0].text = "Problem Statement"
            if len(tf.paragraphs) > 1:
                tf.paragraphs[1].text = "Challenges in Statutory Schedule III Financial Ratio Reporting"
        elif shape.name == "TextBox 8":
            tf = shape.text_frame
            tf.paragraphs[0].text = "Mandatory MCA Compliance Burden"
            tf.paragraphs[1].text = "MCA notification (G.S.R. 207(E)) mandates disclosure of 11 statutory analytical ratios in Schedule III Balance Sheet notes."
            tf.paragraphs[2].text = "• Manual computation across CY & PY opening/closing averages is time-intensive and highly prone to formula errors."
            tf.paragraphs[3].text = "• Rigorous ICAI Guidance Note formula definitions and statutory zero-handling must be strictly followed."
        elif shape.name == "TextBox 11":
            tf = shape.text_frame
            tf.paragraphs[0].text = "Non-Standard Financial Ingestion"
            tf.paragraphs[1].text = "Auditors and enterprises receive financial statements in fragmented, non-uniform formats."
            tf.paragraphs[2].text = "• Disparate Excel layouts with merged cells, dynamic header rows, trailing sheet name spaces, and parenthesized negatives."
            tf.paragraphs[3].text = "• 40+ naming variations and splits (e.g. MSME Trade Payables, Fixed Assets vs PPE, Capital Employed)."
        elif shape.name == "TextBox 14":
            tf = shape.text_frame
            tf.paragraphs[0].text = "Variance Analysis & Narrative Gap"
            tf.paragraphs[1].text = "Statutory rule requires mandatory explanation for every ratio varying by 25% or more."
            tf.paragraphs[2].text = "• Identifying underlying numerical drivers across 11 ratios manually creates audit compliance bottlenecks."
            tf.paragraphs[3].text = "• Drafting auditor-grade narrative reasoning naming exact amounts and percentage drivers is repetitive and tedious."

    # --------------------------------------------------------------------------
    # SLIDE 3: Technology Used
    # --------------------------------------------------------------------------
    s3 = prs.slides[2]
    for shape in s3.shapes:
        if shape.name == "TextBox 3":
            tf = shape.text_frame
            tf.paragraphs[0].text = "Technology Used"
            if len(tf.paragraphs) > 1:
                tf.paragraphs[1].text = "Modern, Robust & Secure Desktop Architecture"
        elif shape.name == "TextBox 7":
            tf = shape.text_frame
            tf.paragraphs[0].text = "⚙️ Core Engine & GUI  |  PySide6 & Python"
            tf.paragraphs[1].text = "✔ Python 3.14: Pure mathematical engine with zero hardcoded figures"
            tf.paragraphs[2].text = "✔ PySide6 (Qt 6): High-DPI native desktop GUI with tabbed navigation"
            tf.paragraphs[3].text = "✔ SQLite: Embedded relational database for client portfolios"
        elif shape.name == "TextBox 9":
            tf = shape.text_frame
            tf.paragraphs[0].text = "📄 Ingestion & NLP  |  Dynamic Excel Parsing"
            tf.paragraphs[1].text = "✔ openpyxl: Dynamic header scanning & parenthesized negative parser"
            tf.paragraphs[2].text = "✔ RapidFuzz: Accounting synonym mapping across 40+ statutory heads"
            tf.paragraphs[3].text = "✔ Deterministic Rules 1–6: Autonomous ambiguity & duplicate resolution"
        elif shape.name == "TextBox 11":
            tf = shape.text_frame
            tf.paragraphs[0].text = "💻 Audit Workspace  |  Interactive Controls"
            tf.paragraphs[1].text = "✔ Dynamic Threshold Slider: Live recalculation of flagged variances"
            tf.paragraphs[2].text = "✔ Statutory Guidance Popups (ℹ️): Schedule III clauses & audit insights"
            tf.paragraphs[3].text = "✔ Multiline Reason Editor: Interactive auditor discretion override"
        elif shape.name == "TextBox 13":
            tf = shape.text_frame
            tf.paragraphs[0].text = "📊 Export & Distribution  |  Packaging"
            tf.paragraphs[1].text = "✔ python-docx: Audit-ready A4 Landscape Word reports with tables"
            tf.paragraphs[2].text = "✔ openpyxl: Live formula-driven Schedule III Excel workpapers"
            tf.paragraphs[3].text = "✔ PyInstaller: Standalone single-file Windows executable (.exe)"
        elif shape.name == "Rounded Rectangle 14":
            tf = shape.text_frame
            tf.paragraphs[0].text = "🔒 Data Storage & Security: All client records and uploaded financial data are stored 100% locally on the user's workstation in an embedded SQLite database at %APPDATA%\\ScheduleIIIRatioAnalyser\\ratio_analyser.db. 100% offline air-gapped execution with zero cloud storage, zero external telemetry, and SHA-256 data integrity verification, guaranteeing complete client financial confidentiality."

    # --------------------------------------------------------------------------
    # SLIDE 4: Implementation
    # --------------------------------------------------------------------------
    s4 = prs.slides[3]
    for shape in s4.shapes:
        if shape.name == "TextBox 3":
            tf = shape.text_frame
            tf.paragraphs[0].text = "Implementation"
            if len(tf.paragraphs) > 1:
                tf.paragraphs[1].text = "Step-by-Step Workflow & Key Solution Capabilities"
        elif shape.name == "TextBox 8":
            tf = shape.text_frame
            tf.paragraphs[0].text = "Zero-Intervention Ingestion"
            tf.paragraphs[1].text = "• 3-Action Workflow: Enter Client Name → Upload CY File → Upload PY File."
            tf.paragraphs[2].text = "• Automated dynamic parsing with sheet/column alias resolution."
            tf.paragraphs[3].text = "• Storage: Processed 100% locally with zero cloud transmission."
        elif shape.name == "TextBox 11":
            tf = shape.text_frame
            tf.paragraphs[0].text = "11 Mandated Ratios Engine"
            tf.paragraphs[1].text = "• Computes all 11 MCA Schedule III ratios strictly per ICAI norms."
            tf.paragraphs[2].text = "• Automated 2-year opening/closing average balance handling."
            tf.paragraphs[3].text = "• 3-step DSCR principal repayment waterfall and divide-by-zero protection."
        elif shape.name == "TextBox 14":
            tf = shape.text_frame
            tf.paragraphs[0].text = "Smart Variance Analyzer"
            tf.paragraphs[1].text = "• Dynamically flags ratios exceeding statutory 25% threshold."
            tf.paragraphs[2].text = "• Decomposes underlying drivers into natural narrative explanations."
            tf.paragraphs[3].text = "• Allows auditor free-text edits with persistent database storage."
        elif shape.name == "TextBox 17":
            tf = shape.text_frame
            tf.paragraphs[0].text = "Statutory Integrity & Exports"
            tf.paragraphs[1].text = "• Automated IC-1 to IC-9 cross-statement consistency checks."
            tf.paragraphs[2].text = "• 1-click export to Board-ready Word (.docx) & formula Excel (.xlsx)."
            tf.paragraphs[3].text = "• Includes statutory clauses, footnotes, and mathematical drilldown."

    # --------------------------------------------------------------------------
    # SLIDE 5: Conclusion
    # --------------------------------------------------------------------------
    s5 = prs.slides[4]
    for shape in s5.shapes:
        if shape.name == "TextBox 3":
            tf = shape.text_frame
            tf.paragraphs[0].text = "Conclusion"
            if len(tf.paragraphs) > 1:
                tf.paragraphs[1].text = "Project Impact, Audit Utility & Enterprise Delivery"
        elif shape.name == "Rounded Rectangle 6":
            tf = shape.text_frame
            tf.paragraphs[0].text = "PROJECT CONCLUSION"
            tf.paragraphs[1].text = "Transforming Statutory Audit Reporting"
            tf.paragraphs[2].text = "⚡ 90%+ Time Savings: Automates statutory ratio computations, reviews, and working paper generation in under 5 seconds."
            tf.paragraphs[3].text = "🎯 100% Regulatory Compliance: Complete alignment with MCA Schedule III (2021) & ICAI Guidance Notes."
            tf.paragraphs[4].text = "🔒 Data Privacy Assured: 100% local SQLite desktop architecture with zero cloud storage footprint."
            tf.paragraphs[5].text = "💼 Audit Value Creation: Shifts CA bandwidth from mechanical data re-entry to high-level advisory judgment."
        elif shape.name == "Rounded Rectangle 7":
            tf = shape.text_frame
            tf.paragraphs[0].text = "CORE CONCLUSION"
            tf.paragraphs[1].text = "Audit-Grade Utility for Chartered Accountants"
            tf.paragraphs[2].text = "Successfully bridges multi-format document friction, eliminates computational error risks, performs IC-1 to IC-9 integrity checks, and automatically drafts statutory 25% variance explanations ready for Board and audit working papers."
        elif shape.name == "Rounded Rectangle 8":
            tf = shape.text_frame
            tf.paragraphs[0].text = "STANDALONE DEPLOYMENT"
            tf.paragraphs[1].text = "📦 Production Executable (.exe) & Test Automation"
            tf.paragraphs[2].text = "1. Single-File Windows Executable: Packaged as ScheduleIIIRatioAnalyser.exe for instant zero-dependency deployment."
            tf.paragraphs[3].text = "• 100% Automated Test Suite: 34 statutory acceptance and engine unit tests passing with zero hardcoded figures."

    # Save to both locations
    prs.save(pptx_source)
    dest_path = r"C:\Users\7PIN\Desktop\ICAI_AI_Course\Level_2\Capstone\Schedule_III_Ratio_Analyser_Presentation.pptx"
    prs.save(dest_path)
    print(f"Updated successfully:\n1. {pptx_source}\n2. {dest_path}")

if __name__ == "__main__":
    update_presentation()

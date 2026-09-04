"""Generate a professional 5-slide Capstone Project presentation in PowerPoint (.pptx) format."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Corporate Color Palette
NAVY = RGBColor(10, 37, 64)       # #0A2540
BLUE = RGBColor(0, 102, 204)      # #0066CC
SKY_BG = RGBColor(240, 246, 255)  # #F0F6FF
CARD_BG = RGBColor(248, 250, 252) # #F8FAFC
BORDER = RGBColor(226, 232, 240)  # #E2E8F0
TEXT_DARK = RGBColor(15, 23, 42)  # #0F172A
TEXT_MUTED = RGBColor(100, 116, 139) # #64748B
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(15, 118, 110)     # #0F766E


def create_deck(output_path: str):
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Helper: Add Header Bar to standard slides
    def add_slide_header(slide, title_text, category="CAPSTONE PROJECT"):
        # Top banner shape
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        banner.fill.solid()
        banner.fill.fore_color.rgb = NAVY
        banner.line.color.rgb = NAVY

        # Category / Tag
        tx_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.18), Inches(10), Inches(0.3))
        tf_cat = tx_cat.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = RGBColor(186, 230, 253) # Light Blue

        # Title
        tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.55))
        tf_title = tx_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE

    # Helper: Add Card Box
    def add_card(slide, left, top, width, height, title, items, bg_color=CARD_BG, border_color=BORDER, accent_color=BLUE):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)

        # Content Box
        tx = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = tx.text_frame
        tf.word_wrap = True
        
        # Card Title
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = accent_color
        p_t.space_after = Pt(10)

        # Items
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(12.5)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    
    # Background Hero
    hero = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()

    # Subtle decorative accent bar
    accent = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.3), Inches(0.15), Inches(4.8))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()

    # Title content
    tb1 = s1.shapes.add_textbox(Inches(1.4), Inches(1.5), Inches(10.5), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p_badge = tf1.paragraphs[0]
    p_badge.text = "ICAI AI COURSE • LEVEL 2 CAPSTONE PROJECT"
    p_badge.font.size = Pt(13)
    p_badge.font.bold = True
    p_badge.font.color.rgb = RGBColor(56, 189, 248) # Cyan Accent
    p_badge.space_after = Pt(14)

    p_main = tf1.add_paragraph()
    p_main.text = "Schedule III Ratio Analyser"
    p_main.font.size = Pt(36)
    p_main.font.bold = True
    p_main.font.color.rgb = WHITE
    p_main.space_after = Pt(8)

    p_sub = tf1.add_paragraph()
    p_sub.text = "Autonomous Statutory Analytical Ratios & Variance Intelligence Engine"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = RGBColor(226, 232, 240)
    p_sub.space_after = Pt(28)

    p_desc = tf1.add_paragraph()
    p_desc.text = "An offline desktop application for Chartered Accountants and statutory audit teams to compute, validate, and disclose the 11 Analytical Ratios mandated under Schedule III of the Companies Act, 2013."
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = RGBColor(148, 163, 184)
    p_desc.space_after = Pt(32)

    p_author = tf1.add_paragraph()
    p_author.text = "Submitted By: xxx"
    p_author.font.size = Pt(15)
    p_author.font.bold = True
    p_author.font.color.rgb = RGBColor(254, 240, 138) # Gold accent

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_header(s2, "Problem Statement: Challenges in Statutory Analytical Ratios")

    add_card(
        s2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
        "📜 The Regulatory Mandate",
        [
            "MCA Notification G.S.R. 207(E) amended Schedule III to mandate 11 statutory analytical ratios in audited financial statements.",
            "Mandatory Disclosure: Explanation is statutorily required for any ratio variance exceeding 25% compared to the preceding year.",
            "Rigorous Statutory Standard: Strict adherence required for numerator/denominator definitions, 2-year averages, and divide-by-zero handling."
        ],
        accent_color=NAVY
    )

    add_card(
        s2, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4),
        "⚠️ Key Audit & Industry Bottlenecks",
        [
            "Manual & Error-Prone: Calculating 11 ratios, 2-year averages, and driver decompositions across dozens of audit clients is tedious and prone to formula errors.",
            "Grouping & Label Discrepancies: Client Excel workbooks use inconsistent labels (e.g. Fixed Assets vs PPE, MSME dues splits, trailing spaces in sheet names).",
            "Complex Debt Reconciliation: Cash flow financing flows often do not articulate with balance sheet debt movements.",
            "Repetitive Drafting: Drafting professional statutory reasons naming exact underlying drivers consumes valuable audit hours."
        ],
        accent_color=RGBColor(192, 57, 43)
    )

    # -------------------------------------------------------------
    # SLIDE 3: Technology Used & Data Security
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_header(s3, "Technology Stack, Data Storage & Enterprise Security")

    add_card(
        s3, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
        "💻 Modern Technology Stack",
        [
            "Python 3.11+: Pure mathematical calculation engine with zero hardcoded figures.",
            "PySide6 (Qt 6): High-DPI native desktop GUI with tabbed navigation and interactive dialogs.",
            "openpyxl & python-docx: Dynamic Excel ingestion and audit-ready A4 landscape Word export generation.",
            "RapidFuzz: High-performance fuzzy aliasing for intelligent statement line-item normalization.",
            "PyInstaller: Standalone single-file Windows executable (.exe) for zero-setup distribution."
        ],
        accent_color=BLUE
    )

    add_card(
        s3, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4),
        "🔒 Local Data Storage & Complete Security",
        [
            "Where Data is Stored: 100% locally on the auditor's workstation inside an embedded SQLite database at %APPDATA%\\ScheduleIIIRatioAnalyser\\ratio_analyser.db.",
            "100% Offline & Air-Gapped: Zero network requests, zero cloud API dependencies, and zero telemetry.",
            "Strict Client Confidentiality: Client financial statements never leave the local computer.",
            "Cryptographic Audit Trail: SHA-256 cryptographic hashes computed for all ingested workbooks to guarantee data integrity.",
            "Database Backup & Restore: Built-in 1-click database backup and restore for disaster recovery."
        ],
        accent_color=TEAL
    )

    # -------------------------------------------------------------
    # SLIDE 4: Implementation
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_header(s4, "Implementation: Architecture & Autonomous Workflow")

    add_card(
        s4, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.4),
        "1️⃣ Zero-Intervention Pipeline",
        [
            "3-Action Workflow: Enter Client Name → Upload CY File → Upload PY File.",
            "Autonomous Execution: Complete analysis runs unattended to Screen 4.",
            "No Blocking Prompts: Deterministic Rules 1–6 resolve duplicate captions, MSME splits, and aliases automatically."
        ],
        accent_color=NAVY
    )

    add_card(
        s4, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.4),
        "2️⃣ Core Analytical Engine",
        [
            "11 Schedule III Ratios: Full mathematical coverage with 2-year opening/closing averages.",
            "Dynamic Waterfall: 3-step debt service waterfall for DSCR principal repayments.",
            "Variance Engine: Automatically generates multi-line statutory reasons naming actual amounts and % drivers."
        ],
        accent_color=BLUE
    )

    add_card(
        s4, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.4),
        "3️⃣ Audit Features & Exports",
        [
            "Statutory Integrity (IC-1 to IC-9): Dynamic cross-statement & cash flow consistency checks.",
            "Ratio Guidance (ℹ️): Instant popups with statutory clauses & audit significance.",
            "1-Click Word & Excel Exports: Audit-ready A4 landscape reports ready for filing."
        ],
        accent_color=TEAL
    )

    # -------------------------------------------------------------
    # SLIDE 5: Conclusion
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_header(s5, "Conclusion: Key Value Delivered & Impact")

    add_card(
        s5, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
        "⭐ Key Business Outcomes",
        [
            "90%+ Time Savings: Reduces hours of manual ratio computations and drafting to under 5 seconds.",
            "100% Regulatory Accuracy: Eliminates calculation discrepancies and guarantees Schedule III formula compliance.",
            "Standardized Audit Quality: Delivers uniform, high-quality statutory notes across all clients and engagement teams.",
            "Professional Reporting: Exports clean, formatted Word (.docx) and Excel (.xlsx) files directly into audit working papers."
        ],
        accent_color=BLUE
    )

    add_card(
        s5, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4),
        "🚀 Project Highlights & Delivery",
        [
            "100% Test Automation: 34 automated unit and statutory acceptance tests passing with 100% success rate.",
            "Zero Hardcoding Constraint: Engine dynamically parses all financial structures without hardcoded figures.",
            "Standalone Executable: Single-file Windows binary (ScheduleIIIRatioAnalyser.exe) ready for instant enterprise deployment.",
            "Future-Ready: Fully extensible architecture supporting additional industry-specific ratios and custom templates."
        ],
        accent_color=TEAL
    )

    prs.save(output_path)
    print(f"Presentation successfully saved to: {output_path}")


if __name__ == "__main__":
    out_file = "Schedule_III_Ratio_Analyser_Presentation.pptx"
    create_deck(out_file)

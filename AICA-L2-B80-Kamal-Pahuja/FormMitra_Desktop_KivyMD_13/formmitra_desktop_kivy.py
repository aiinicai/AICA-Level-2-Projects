#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
FormMitra Desktop (Kivy/KivyMD edition) - Assistive Form-Reading App
========================================================================

This is a UI/UX-focused rewrite of formmitra_desktop.py (the original
Tkinter version). Same purpose, same features, same underlying logic
(Gemini calls, form loaders, font handling, saved-copy rendering,
text-to-speech) - the whole application layer has been rebuilt on
KivyMD (Material Design components on top of Kivy) instead of Tkinter,
for a more modern look: proper cards, themed buttons, smooth dialogs,
light/dark theming, and an embedded live camera preview instead of a
separate popup window.

Purpose
-------
Helps elderly users, or anyone who finds official paperwork confusing,
understand blank forms (bank account opening, insurance, demat, KYC, etc.)
in a regional language of their choice.

How to run
----------
1. Double-click this file to open it in IDLE, then press F5 (Run Module).
   OR run from a terminal:  python formmitra_desktop_kivy.py
2. The FIRST time you run it, the app will automatically check for and
   install every third-party library it needs (Kivy, KivyMD, Pillow,
   OpenCV, PyMuPDF, python-docx, openpyxl, python-pptx, google-genai,
   gTTS, playsound). This can take a few minutes and needs an internet
   connection - Kivy/KivyMD in particular are larger downloads, so the
   very first run may take longer than the Tkinter version did. You will
   see progress messages in the IDLE Shell / terminal window.
3. On first use you will be asked for a free Google AI Studio (Gemini)
   API key. Get one at https://aistudio.google.com/app/apikey
   The key is saved locally on your own computer only (in a file under
   your home folder) and is never sent anywhere except directly to
   Google's Gemini API.
4. Tap "Upload" (image / PDF / Word / Excel / PowerPoint) or "Camera" to
   capture a blank form, pick your language, and read the bright guidance
   notes that appear on the form. Tap any numbered badge to reveal just
   that field's note.

This is a single, self-contained file on purpose, so it is easy to run
directly from IDLE without installing a package.
"""

import sys
import os
import time
import subprocess
import importlib

# ---------------------------------------------------------------------------
# STEP 0 — SELF-INSTALLING BOOTSTRAP
# ---------------------------------------------------------------------------
# Everything in this section uses ONLY the Python standard library, because
# it runs before we know whether any third-party package is available yet.

# Maps: the name we `import` -> the name pip installs it under
REQUIRED_PACKAGES = {
    "kivy": "kivy>=2.3.1",
    "kivymd": "kivymd==1.2.0",
    "PIL": "Pillow",
    "cv2": "opencv-python",
    "fitz": "PyMuPDF",
    "docx": "python-docx",
    "openpyxl": "openpyxl",
    "pptx": "python-pptx",
    "google.genai": "google-genai",
    "gtts": "gTTS",
    "playsound": "playsound==1.2.2",
    "requests": "requests",
}


def _run_pip(cmd):
    """Run a pip command, capturing output so failures are actually visible
    (plain subprocess.check_call with --quiet hides pip's real error, which
    just leaves the user staring at a confusing crash later on)."""
    result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True)
    return result.returncode == 0, result.stdout


def _pip_install(pip_name):
    """Install one package with pip, trying a normal install first, then a
    user-site install. Prints pip's actual output on failure so the real
    cause (no wheel for this Python version, missing compiler, network
    block, permissions, ...) is visible instead of a bare error code."""
    print(f"  Installing {pip_name} ... (this can take a minute)")
    cmd_base = [sys.executable, "-m", "pip", "install", "--upgrade"]
    ok, output = _run_pip(cmd_base + [pip_name])
    if ok:
        return True
    ok, output2 = _run_pip(cmd_base + ["--user", pip_name])
    if ok:
        return True
    print(f"  !! Could not install {pip_name}. pip's own error output:")
    tail = output2.strip().splitlines() or output.strip().splitlines()
    for line in tail[-15:]:
        print(f"       {line}")
    return False


def ensure_dependencies():
    """Check every required library; auto-install anything missing."""
    print("FormMitra: checking required libraries ...")
    missing = []
    for import_name in REQUIRED_PACKAGES:
        try:
            importlib.import_module(import_name)
        except ImportError:
            missing.append(import_name)

    if not missing:
        print("FormMitra: all required libraries are already installed.\n")
        return

    print(f"FormMitra: {len(missing)} librar{'y is' if len(missing)==1 else 'ies are'} "
          f"missing and will be installed automatically:")
    for name in missing:
        print(f"   - {name}  (pip package: {REQUIRED_PACKAGES[name]})")
    print("  (Kivy/KivyMD are larger downloads - the first run can take a few minutes.)")
    print()

    # Make sure pip itself is available.
    try:
        importlib.import_module("pip")
    except ImportError:
        print("FormMitra: pip is not available; trying to bootstrap it with ensurepip ...")
        try:
            subprocess.check_call([sys.executable, "-m", "ensurepip", "--upgrade"])
        except Exception as exc:
            print("FormMitra: could not bootstrap pip automatically:", exc)
            print("Please install pip manually, then re-run this program.")
            sys.exit(1)

    failed = []
    for import_name in missing:
        pip_name = REQUIRED_PACKAGES[import_name]
        ok = _pip_install(pip_name)
        if ok:
            try:
                importlib.invalidate_caches()
                importlib.import_module(import_name)
                print(f"  OK: {import_name} is ready.")
            except ImportError:
                failed.append(import_name)
        else:
            failed.append(import_name)

    if failed:
        print("\nFormMitra: the following libraries could not be installed "
              "automatically:")
        for name in failed:
            print(f"   - {name} (pip install {REQUIRED_PACKAGES[name]})")

        # Kivy and KivyMD build the app's entire window/UI - there is no
        # partial mode without them, so don't limp forward into a confusing
        # crash. Everything else (PDF support, TTS, etc.) can degrade
        # gracefully at runtime instead.
        core_ui = {"kivy", "kivymd"}
        if core_ui & set(failed):
            print("\nFormMitra can't run without Kivy/KivyMD (they build the app's "
                  "window). Please open a Command Prompt / terminal in this folder "
                  "and run, one at a time:")
            for name in sorted(core_ui & set(failed)):
                print(f"   {sys.executable} -m pip install --upgrade {REQUIRED_PACKAGES[name]}")
            print("\nRunning it from a terminal (instead of double-clicking, or F5 in "
                  "IDLE) lets you see pip's real error message above, which is usually "
                  "one of: no matching wheel for your Python version (try upgrading pip "
                  "first: python -m pip install --upgrade pip), a blocked network/proxy, "
                  "or a permissions issue.")

            if sys.version_info >= (3, 14):
                print(f"\nYou're running Python {sys.version_info.major}."
                      f"{sys.version_info.minor} — as of today, Kivy's Windows build "
                      "dependencies (specifically kivy_deps.sdl2_dev) have not been "
                      "updated with a Python 3.14 build yet, so Kivy currently CANNOT "
                      "be installed on Python 3.14 on Windows at all, no matter what "
                      "you try here. This is a gap in Kivy's own release, not something "
                      "this script can work around.\n"
                      "Fix: install Python 3.12 alongside your current Python (this "
                      "does not remove or affect 3.14) - if you installed Python via "
                      "the newer python.org/Microsoft Store installer, an easy way is "
                      "opening a terminal and running:  py install 3.12\n"
                      "Otherwise, download the Python 3.12 installer from "
                      "python.org/downloads and run it (leave the box for 'tcl/tk and "
                      "IDLE' checked).\n"
                      "Then run this script specifically with that version:\n"
                      "   py -3.12 " + os.path.basename(sys.argv[0] or "formmitra_desktop_kivy.py") + "\n"
                      "(or open THIS folder, right-click formmitra_desktop_kivy.py -> "
                      "Open with -> choose the Python 3.12 you just installed)\n"
                      "If you'd rather not install another Python version right now, "
                      "the Tkinter edition (formmitra_desktop.py) has no such Kivy "
                      "dependency and will run fine on Python 3.14 as-is.\n")
            elif sys.version_info < (3, 8):
                print(f"\nYou're running Python {sys.version_info.major}."
                      f"{sys.version_info.minor}, which is older than Kivy currently "
                      "supports. Installing a newer Python (3.11 or 3.12 are safe, "
                      "well-supported choices) from python.org should fix this.\n")
            else:
                print("\nIf the error above mentions your Python version specifically "
                      "(e.g. 'no matching distribution'), installing Python 3.11 or "
                      "3.12 from python.org alongside your current version and running "
                      "this script with that one usually fixes it, since Kivy's wheels "
                      "lag behind the newest Python releases.\n")
            sys.exit(1)

        print("\nSome features may not work until the remaining libraries above are "
              "installed by hand, for example by running that pip command yourself "
              "in a terminal. The app will still try to start.\n")
    else:
        print("\nFormMitra: all libraries installed successfully.\n")


ensure_dependencies()

# ---------------------------------------------------------------------------
# STEP 1 — NOW SAFE TO IMPORT EVERYTHING
# ---------------------------------------------------------------------------
import json
import re
import queue
import threading
import tempfile
import traceback

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    Image = None

try:
    import cv2  # noqa: F401  (used by Kivy's own camera provider on desktop)
except ImportError:
    cv2 = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from google import genai
    from google.genai import types as genai_types
except ImportError:
    genai = None
    genai_types = None

try:
    from gtts import gTTS
except ImportError:
    gTTS = None

try:
    import playsound
except ImportError:
    playsound = None

from kivy.config import Config
Config.set("graphics", "width", "1280")
Config.set("graphics", "height", "820")

from kivy.app import App
from kivy.clock import Clock
from kivy.metrics import dp, sp
from kivy.uix.floatlayout import FloatLayout
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.widget import Widget
from kivy.uix.scrollview import ScrollView
from kivy.uix.gridlayout import GridLayout
from kivy.uix.image import Image as KivyImage
from kivy.uix.spinner import Spinner
from kivy.core.text import LabelBase
from kivy.graphics import Color, Ellipse, Rectangle, RoundedRectangle, Line
from kivy.graphics.texture import Texture

from kivymd.app import MDApp
from kivymd.uix.boxlayout import MDBoxLayout
from kivymd.uix.card import MDCard
from kivymd.uix.label import MDLabel
from kivymd.uix.button import (
    MDRaisedButton, MDIconButton, MDFlatButton,
    MDFillRoundFlatIconButton, MDRoundFlatIconButton,
)
from kivymd.uix.textfield import MDTextField
from kivymd.uix.selectioncontrol import MDSwitch
from kivymd.uix.dialog import MDDialog
from kivymd.uix.snackbar import Snackbar
from kivymd.uix.toolbar import MDTopAppBar
from kivymd.uix.menu import MDDropdownMenu
from kivymd.uix.list import OneLineIconListItem, OneLineAvatarIconListItem, IconLeftWidget


# ---------------------------------------------------------------------------
# CONSTANTS
# ---------------------------------------------------------------------------

APP_NAME = "FormMitra"
CONFIG_DIR = os.path.join(os.path.expanduser("~"), ".formmitra")
CONFIG_PATH = os.path.join(CONFIG_DIR, "config.json")
FONT_CACHE_DIR = os.path.join(CONFIG_DIR, "fonts")

# Candidate Gemini models, tried in order (first that works is used & cached).
# NOTE: Google renames/retires Gemini model IDs fairly often. This list is
# tried top-to-bottom and the first that works is cached for the rest of
# the session; GeminiClient also falls back to auto-discovering a working
# model via the API's own ListModels if every name below has been retired
# by the time you're reading this - see GeminiClient._discover_model().
GEMINI_MODEL_CANDIDATES = [
    "gemini-3.5-flash",
    "gemini-3.1-flash-lite",
    "gemini-2.5-flash",
    "gemini-2.5-flash-lite",
]

LANGUAGES = [
    "English", "Hindi", "Bengali", "Tamil", "Telugu", "Marathi",
    "Gujarati", "Kannada", "Malayalam", "Punjabi", "Odia",
    "Urdu", "Assamese",
]

# Rough script family per language, used for font resolution.
LANGUAGE_SCRIPT = {
    "Hindi": "Devanagari", "Marathi": "Devanagari",
    "Bengali": "Bengali", "Assamese": "Bengali",
    "Tamil": "Tamil", "Telugu": "Telugu", "Kannada": "Kannada",
    "Malayalam": "Malayalam", "Gujarati": "Gujarati",
    "Punjabi": "Gurmukhi", "Odia": "Odia", "Urdu": "Arabic",
    "English": "Latin",
}

# Script -> Noto Sans family name, for downloading/searching a font file
# that actually contains that script's glyphs. Verified working source:
# the notofonts.github.io mirror, which serves plain static .ttf files
# (unlike the main google/fonts repo, which now only ships variable fonts
# under less predictable paths).
NOTO_FONT_FAMILY = {
    "Devanagari": "NotoSansDevanagari",
    "Bengali": "NotoSansBengali",
    "Tamil": "NotoSansTamil",
    "Telugu": "NotoSansTelugu",
    "Kannada": "NotoSansKannada",
    "Malayalam": "NotoSansMalayalam",
    "Gujarati": "NotoSansGujarati",
    "Gurmukhi": "NotoSansGurmukhi",
    "Odia": "NotoSansOriya",     # Noto still uses the older "Oriya" name
    "Arabic": "NotoSansArabic",  # used for Urdu
}

# gTTS language codes for the "read aloud" feature (best effort).
TTS_LANG_CODES = {
    "English": "en", "Hindi": "hi", "Bengali": "bn", "Tamil": "ta",
    "Telugu": "te", "Marathi": "mr", "Gujarati": "gu", "Kannada": "kn",
    "Malayalam": "ml", "Punjabi": "pa", "Odia": "or", "Urdu": "ur",
    "Assamese": "as",
}

# Native digit glyphs per script, and short fixed UI words ("Required" /
# "Example") that appear next to every field's guidance. Both the saved
# copy (Pillow, baking text into a picture) and this app's on-screen
# guidance render a whole string with ONE font file and no OS-level
# font-linking/fallback the way a browser or Windows' own UI has - so if
# that font is a script-specific one, any literal ASCII digits/words our
# own template code mixes into the string (field numbering, "Required",
# "Example:") can show up as tofu boxes even though the actual translated
# guidance text next to them renders fine. Using each language's own
# digits and its own word for "Required"/"Example" avoids ever needing
# Latin glyphs from that font for text this app controls.
LOCAL_DIGITS = {
    "Devanagari": "०१२३४५६७८९",
    "Bengali": "০১২৩৪৫৬৭৮৯",
    "Tamil": "௦௧௨௩௪௫௬௭௮௯",
    "Telugu": "౦౧౨౩౪౫౬౭౮౯",
    "Kannada": "೦೧೨೩೪೫೬೭೮೯",
    "Malayalam": "൦൧൨൩൪൫൬൭൮൯",
    "Gujarati": "૦૧૨૩૪૫૬૭૮૯",
    "Gurmukhi": "੦੧੨੩੪੫੬੭੮੯",
    "Odia": "୦୧୨୩୪୫୬୭୮୯",
    "Arabic": "۰۱۲۳۴۵۶۷۸۹",  # used for Urdu
}

REQUIRED_WORD = {
    "English": "Required", "Hindi": "आवश्यक", "Marathi": "आवश्यक",
    "Bengali": "প্রয়োজনীয়", "Assamese": "প্ৰয়োজনীয়",
    "Tamil": "கட்டாயம்", "Telugu": "అవసరం", "Kannada": "ಅಗತ್ಯ",
    "Malayalam": "ആവശ്യമാണ്", "Gujarati": "જરૂરી",
    "Punjabi": "ਲੋੜੀਂਦਾ", "Odia": "ଆବଶ୍ୟକ", "Urdu": "ضروری",
}

EXAMPLE_WORD = {
    "English": "Example", "Hindi": "उदाहरण", "Marathi": "उदाहरण",
    "Bengali": "উদাহরণ", "Assamese": "উদাহৰণ",
    "Tamil": "எடுத்துக்காட்டு", "Telugu": "ఉదాహరణ", "Kannada": "ಉದಾಹರಣೆ",
    "Malayalam": "ഉദാഹരണം", "Gujarati": "ઉદાહરણ",
    "Punjabi": "ਉਦਾਹਰਣ", "Odia": "ଉଦାହରଣ", "Urdu": "مثال",
}


def local_number(n, language):
    """Render a field's 1-based index using the target language's own
    digits (e.g. 3 -> "૩" for Gujarati) instead of ASCII "3", so field
    numbering never depends on a script-specific font also including
    Latin/ASCII digit glyphs."""
    script = LANGUAGE_SCRIPT.get(language, "Latin")
    digits = LOCAL_DIGITS.get(script)
    if not digits:
        return str(n)
    return "".join(digits[int(ch)] if ch.isdigit() else ch for ch in str(n))


_TRAILING_ASCII_PUNCT = re.compile(r"[\s.,:;!?]+$")


def _strip_trailing_ascii_punct(text):
    """Drop a trailing ASCII sentence-ending mark (usually just ".") from
    Gemini-translated text for non-Latin languages. Same root issue as
    local_number() above: a stray Latin punctuation character at the very
    end of an otherwise-native-script sentence can render as a tofu box
    if the script font doesn't happen to include it."""
    if not text:
        return text
    return _TRAILING_ASCII_PUNCT.sub("", text)

BRIGHT_NOTE_BG = (255, 235, 59)       # bright yellow (PIL saved-copy colors)
BRIGHT_NOTE_BORDER = (191, 54, 12)
BRIGHT_NOTE_TEXT = (17, 17, 17)
BADGE_BG = (211, 47, 47)
BADGE_TEXT = (255, 255, 255)

BADGE_RGBA = (0.83, 0.18, 0.18, 1)       # same red, as Kivy 0-1 floats
BADGE_ACTIVE_RGBA = (1, 0.43, 0, 1)      # orange, when a note is showing
NOTE_BG_RGBA = (1, 0.92, 0.23, 1)        # bright yellow
NOTE_BORDER_RGBA = (0.75, 0.21, 0.05, 1)

# App-wide accent palette (Kivy 0-1 RGBA). A single teal accent used
# consistently across the top bar, primary buttons, and highlights, instead
# of the flat default-blue Material look.
ACCENT_TEAL = (0.0, 0.47, 0.45, 1)       # primary brand color (buttons, top bar)
ACCENT_TEAL_DARK = (0.0, 0.32, 0.31, 1)  # pressed/darker variant
ACCENT_BLUE = (0.10, 0.46, 0.82, 1)      # secondary accent (Listen / info actions)
SURFACE_TINT = (0.94, 0.97, 0.97, 1)     # faint teal-tinted page background

os.makedirs(CONFIG_DIR, exist_ok=True)
os.makedirs(FONT_CACHE_DIR, exist_ok=True)


# ---------------------------------------------------------------------------
# CONFIG (API key + last-used preferences), stored locally only
# ---------------------------------------------------------------------------

def load_config():
    if os.path.exists(CONFIG_PATH):
        try:
            with open(CONFIG_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {}


def save_config(cfg):
    try:
        with open(CONFIG_PATH, "w", encoding="utf-8") as f:
            json.dump(cfg, f, indent=2)
    except Exception as exc:
        print("FormMitra: could not save settings:", exc)


# ---------------------------------------------------------------------------
# FONT RESOLUTION
# ---------------------------------------------------------------------------
# Unlike the Tkinter version (which needed two DIFFERENT approaches - an
# OS font-family NAME for on-screen Tk widgets, vs an actual font FILE for
# baking text into the saved-copy image with Pillow), Kivy consistently
# uses real font FILES everywhere (Label.font_name points at a registered
# file, not an OS family name). So here there is only ONE font-resolution
# path, used for both on-screen Kivy labels and the saved-copy image - a
# genuine reliability improvement, not just a UI reskin.
# ---------------------------------------------------------------------------

_WINDOWS_FONT_DIR = r"C:\Windows\Fonts"
_NIRMALA_REGULAR = os.path.join(_WINDOWS_FONT_DIR, "Nirmala.ttf")
_NIRMALA_BOLD = os.path.join(_WINDOWS_FONT_DIR, "NirmalaB.ttf")

_PAN_SCRIPT_CANDIDATES = [
    "/Library/Fonts/Arial Unicode.ttf",
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
]

_LATIN_ONLY_CANDIDATES = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
]

_LINUX_NOTO_DIRS = [
    "/usr/share/fonts/truetype/noto",
    "/usr/share/fonts/opentype/noto",
    "/usr/share/fonts/truetype/noto-cjk",
    "/usr/local/share/fonts",
]

_registered_kivy_fonts = set()


def _find_local_noto_font(script):
    if not script or script == "Latin":
        return None
    family = NOTO_FONT_FAMILY.get(script, f"NotoSans{script}")
    wanted = family.lower()
    for base in _LINUX_NOTO_DIRS:
        if not os.path.isdir(base):
            continue
        try:
            for fname in os.listdir(base):
                if fname.lower().startswith(wanted) and fname.lower().endswith((".ttf", ".otf")):
                    return os.path.join(base, fname)
        except OSError:
            continue
    return None


def _download_noto_font(script):
    """Best-effort download of a Noto Sans font for a script, cached locally.
    Uses the notofonts.github.io mirror (a stable static-font source)."""
    if not script or script == "Latin":
        return None
    family = NOTO_FONT_FAMILY.get(script, f"NotoSans{script}")
    dest = os.path.join(FONT_CACHE_DIR, f"{family}-Regular.ttf")
    if os.path.exists(dest) and os.path.getsize(dest) > 10000:
        return dest
    urls = [
        f"https://raw.githubusercontent.com/notofonts/notofonts.github.io/main/"
        f"fonts/{family}/hinted/ttf/{family}-Regular.ttf",
        f"https://raw.githubusercontent.com/notofonts/notofonts.github.io/main/"
        f"fonts/{family}/unhinted/ttf/{family}-Regular.ttf",
    ]
    try:
        import requests
    except ImportError:
        return None
    for url in urls:
        try:
            resp = requests.get(url, timeout=8)
            if resp.status_code == 200 and len(resp.content) > 10000:
                with open(dest, "wb") as f:
                    f.write(resp.content)
                return dest
        except Exception:
            continue
    return None


def prefetch_font(language):
    """Kick off a background download of the font for `language`, if
    needed, so it's already cached by the time the user hits Save."""
    script = LANGUAGE_SCRIPT.get(language, "Latin")
    if script == "Latin" or script not in NOTO_FONT_FAMILY:
        return
    if os.name == "nt" and os.path.exists(_NIRMALA_REGULAR):
        return

    def _bg():
        try:
            _download_noto_font(script)
        except Exception:
            pass
    threading.Thread(target=_bg, daemon=True).start()


def resolve_font_path(language, bold=False):
    """Return (font_file_path_or_None, script_ok: bool).

    script_ok is only True when we're reasonably confident the font
    actually has glyphs for the target language's script. A None path
    means "use the app's default font" (fine for English/Latin only).
    """
    script = LANGUAGE_SCRIPT.get(language, "Latin")

    if os.name == "nt":
        nirmala = _NIRMALA_BOLD if bold else _NIRMALA_REGULAR
        if os.path.exists(nirmala):
            return nirmala, True

    for path in _PAN_SCRIPT_CANDIDATES:
        if path and os.path.exists(path):
            return path, True

    if script != "Latin":
        for finder in (_find_local_noto_font, _download_noto_font):
            found = finder(script)
            if found:
                return found, True

    for path in _LATIN_ONLY_CANDIDATES:
        if path and os.path.exists(path):
            return path, (script == "Latin")

    return None, (script == "Latin")


def resolve_font(language, size, bold=False):
    """PIL ImageFont version, for baking text into the saved-copy image."""
    path, script_ok = resolve_font_path(language, bold=bold)
    if path:
        try:
            return ImageFont.truetype(path, size), script_ok
        except Exception:
            pass
    try:
        return ImageFont.load_default(), (LANGUAGE_SCRIPT.get(language, "Latin") == "Latin")
    except Exception:
        return ImageFont.load_default(), False


def kivy_font_name(language, bold=False):
    """Registers (once) and returns a Kivy font name usable as
    Label.font_name for on-screen text in this language. Returns None to
    mean "use Kivy's default font" (fine for English)."""
    path, _ = resolve_font_path(language, bold=bold)
    if not path:
        return None
    reg_name = f"lang_{language}_{'b' if bold else 'r'}"
    if reg_name not in _registered_kivy_fonts:
        try:
            LabelBase.register(name=reg_name, fn_regular=path)
            _registered_kivy_fonts.add(reg_name)
        except Exception:
            return None
    return reg_name


# ---------------------------------------------------------------------------
# GEMINI (Google AI Studio) CLIENT — identical to the Tkinter version
# ---------------------------------------------------------------------------

FIELD_JSON_INSTRUCTIONS = """
You are helping an elderly or first-time user fill out a blank official
form (bank / insurance / demat / KYC / government form). Look at the form
and list every field the user needs to fill in (blanks, boxes, checkboxes,
signature lines, etc).

For EACH field return:
- "label": the exact printed label/heading of the field as it appears on
  the form (keep it in the form's original language/script).
- "explanation_en": a very short, simple, plain-English sentence (max ~18
  words) telling a confused or elderly user what to write here. Avoid
  jargon; explain any abbreviation.
- "example_en": one short realistic example value (max ~8 words), or ""
  if not applicable (e.g. for a signature line).
- "required": true or false - whether the form marks this as mandatory.
- "box_2d": the field's bounding box on the image as [ymin, xmin, ymax,
  xmax], using integers from 0 to 1000 (0,0 = top-left of the image,
  1000,1000 = bottom-right).

Return ONLY a JSON array of objects with exactly these five keys, no other
text, no markdown fences.
"""

TEXT_FIELD_JSON_INSTRUCTIONS = """
You are helping an elderly or first-time user fill out a blank official
form. Below is a numbered list of text lines extracted from the form
(each has an "id"). Decide which lines are fillable FIELDS the user must
respond to (a label such as "Name:", "PAN Number", "Date of Birth",
a checkbox option, a signature line, etc) as opposed to plain instructions,
titles, or decorative text.

For EACH line that IS a field, return an object with:
- "id": the same id number given for that line.
- "label": the field's label, cleaned up.
- "explanation_en": a very short, simple, plain-English sentence (max ~18
  words) telling a confused or elderly user what to write here.
- "example_en": one short realistic example value (max ~8 words), or "".
- "required": true or false - your best guess.

Return ONLY a JSON array of these objects (skip lines that are not
fields), no other text, no markdown fences.

Lines:
{lines}
"""

TRANSLATE_INSTRUCTIONS = """
Translate the "label", "explanation" and "example" of every item below
into {language} (use the native script of that language, not
transliteration). Keep translations short, simple and friendly, suitable
for an elderly or first-time reader. Keep the same order and the same
"id" for each item.

Return ONLY a JSON array of objects, each with exactly the keys
"id", "label", "explanation", "example" - no other text, no markdown
fences.

Items:
{items}
"""


class GeminiError(RuntimeError):
    pass


class GeminiClient:
    """Thin wrapper around Google's current `google-genai` SDK
    (the old `google-generativeai` package is end-of-life)."""

    def __init__(self, api_key):
        if genai is None:
            raise GeminiError(
                "The google-genai library isn't installed. "
                "Restart the app so it can be installed automatically."
            )
        self.api_key = api_key
        self.client = genai.Client(api_key=api_key)
        self._working_model = None

    @staticmethod
    def _extract_json_array(text):
        text = text.strip()
        text = re.sub(r"^```(json)?", "", text.strip())
        text = re.sub(r"```$", "", text.strip())
        start = text.find("[")
        end = text.rfind("]")
        if start == -1 or end == -1:
            raise GeminiError("Gemini did not return a JSON list. Raw reply:\n" + text[:500])
        snippet = text[start:end + 1]
        return json.loads(snippet)

    def _discover_model(self):
        try:
            for m in self.client.models.list():
                name = getattr(m, "name", "") or ""
                name = name.split("/")[-1]
                actions = getattr(m, "supported_actions", None) or []
                if "generateContent" not in actions and actions:
                    continue
                if "flash" in name.lower() or "gemini" in name.lower():
                    return name
        except Exception:
            pass
        return None

    def _generate(self, contents, want_json=True):
        config_kwargs = {"temperature": 0.2}
        if want_json:
            config_kwargs["response_mime_type"] = "application/json"
        config = genai_types.GenerateContentConfig(**config_kwargs)

        model_order = [self._working_model] if self._working_model else []
        model_order += [m for m in GEMINI_MODEL_CANDIDATES if m != self._working_model]

        errors = []
        for model_name in model_order:
            try:
                response = self.client.models.generate_content(
                    model=model_name, contents=contents, config=config,
                )
                text = getattr(response, "text", None)
                if not text:
                    raise GeminiError("Gemini returned an empty response.")
                self._working_model = model_name
                return text
            except Exception as exc:  # noqa: BLE001
                errors.append((model_name, exc))
                continue

        combined = " | ".join(str(e) for _, e in errors)
        if any(code in combined for code in ("401", "403", "API_KEY_INVALID", "PERMISSION_DENIED")):
            raise GeminiError(
                "Gemini rejected the request as unauthorized. Double-check your API key "
                "(Settings button) - it may be invalid, expired, or missing permissions.\n\n"
                f"Details: {combined[:400]}"
            )
        if any(code in combined for code in ("429", "RESOURCE_EXHAUSTED", "quota")):
            raise GeminiError(
                "Gemini reports you've hit your usage quota/rate limit. Wait a bit and try "
                f"again, or check your quota at aistudio.google.com.\n\nDetails: {combined[:400]}"
            )

        discovered = self._discover_model()
        if discovered and discovered not in model_order:
            try:
                response = self.client.models.generate_content(
                    model=discovered, contents=contents, config=config,
                )
                text = getattr(response, "text", None)
                if text:
                    self._working_model = discovered
                    return text
            except Exception as exc:  # noqa: BLE001
                errors.append((discovered, exc))

        detail = "\n".join(f"  - {name}: {exc}" for name, exc in errors)
        raise GeminiError(
            "Gemini request failed on every model tried. Google renames/retires model "
            "names over time, so this app's built-in list may be out of date. Full details "
            f"per model tried:\n{detail}"
        )

    def detect_fields_from_image(self, pil_image):
        image = pil_image.convert("RGB")
        text = self._generate([FIELD_JSON_INSTRUCTIONS, image])
        data = self._extract_json_array(text)
        fields = []
        for i, item in enumerate(data):
            try:
                box = item.get("box_2d") or [0, 0, 40, 200]
                fields.append({
                    "id": i,
                    "label": str(item.get("label", "")).strip(),
                    "explanation_en": str(item.get("explanation_en", "")).strip(),
                    "example_en": str(item.get("example_en", "")).strip(),
                    "required": bool(item.get("required", False)),
                    "box_2d": [int(v) for v in box[:4]],
                })
            except Exception:
                continue
        return fields

    def detect_fields_from_lines(self, lines_with_boxes):
        numbered = "\n".join(f"{i}: {ln['text']}" for i, ln in enumerate(lines_with_boxes))
        prompt = TEXT_FIELD_JSON_INSTRUCTIONS.format(lines=numbered)
        text = self._generate([prompt])
        data = self._extract_json_array(text)
        fields = []
        for item in data:
            try:
                idx = int(item.get("id"))
                if idx < 0 or idx >= len(lines_with_boxes):
                    continue
                box = lines_with_boxes[idx]["box_2d"]
                fields.append({
                    "id": idx,
                    "label": str(item.get("label", "")).strip(),
                    "explanation_en": str(item.get("explanation_en", "")).strip(),
                    "example_en": str(item.get("example_en", "")).strip(),
                    "required": bool(item.get("required", False)),
                    "box_2d": box,
                })
            except Exception:
                continue
        return fields

    def translate_fields(self, fields, language):
        if language == "English" or not fields:
            out = []
            for f in fields:
                out.append({**f, "label_t": f["label"], "explanation_t": f["explanation_en"],
                            "example_t": f["example_en"]})
            return out

        items = [
            {"id": f["id"], "label": f["label"], "explanation": f["explanation_en"],
             "example": f["example_en"]}
            for f in fields
        ]
        prompt = TRANSLATE_INSTRUCTIONS.format(
            language=language, items=json.dumps(items, ensure_ascii=False)
        )
        text = self._generate([prompt])
        data = self._extract_json_array(text)
        by_id = {int(item["id"]): item for item in data if "id" in item}

        out = []
        for f in fields:
            t = by_id.get(f["id"])
            if t:
                out.append({**f,
                            "label_t": _strip_trailing_ascii_punct(t.get("label", f["label"])),
                            "explanation_t": _strip_trailing_ascii_punct(
                                t.get("explanation", f["explanation_en"])),
                            "example_t": _strip_trailing_ascii_punct(
                                t.get("example", f["example_en"]))})
            else:
                out.append({**f, "label_t": f["label"], "explanation_t": f["explanation_en"],
                            "example_t": f["example_en"]})
        return out


# ---------------------------------------------------------------------------
# FORM LOADERS — identical to the Tkinter version
# ---------------------------------------------------------------------------

PAGE_W, PAGE_H = 1700, 2200


def load_image_form(path):
    img = Image.open(path)
    return img.convert("RGB")


def load_pdf_form(path, page_number=0):
    if fitz is None:
        raise RuntimeError("PyMuPDF is not installed; cannot open PDF files.")
    doc = fitz.open(path)
    page_count = doc.page_count
    page_number = max(0, min(page_number, page_count - 1))
    page = doc.load_page(page_number)
    zoom = 2.0
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
    mode = "RGB" if pix.n < 4 else "RGBA"
    img = Image.frombytes(mode, [pix.width, pix.height], pix.samples).convert("RGB")
    doc.close()
    return img, page_count


class _LineLayout:
    def __init__(self, title=""):
        self.img = Image.new("RGB", (PAGE_W, PAGE_H), "white")
        self.draw = ImageDraw.Draw(self.img)
        self.y = 60
        self.lines = []
        font_title, _ = resolve_font("English", 34, bold=True)
        self.font, _ = resolve_font("English", 26)
        if title:
            self.draw.text((50, self.y), title, fill="black", font=font_title)
            self.y += 60
            self.draw.line((50, self.y, PAGE_W - 50, self.y), fill=(180, 180, 180), width=2)
            self.y += 30

    def add_line(self, text, indent=0):
        text = (text or "").strip()
        if not text:
            self.y += 12
            return
        x = 50 + indent
        wrapped = self._wrap(text, PAGE_W - 100 - indent)
        start_y = self.y
        for w_line in wrapped:
            if self.y > PAGE_H - 80:
                break
            self.draw.text((x, self.y), w_line, fill="black", font=self.font)
            self.y += 34
        box = [
            int(1000 * start_y / PAGE_H), int(1000 * x / PAGE_W),
            int(1000 * self.y / PAGE_H), int(1000 * (PAGE_W - 50) / PAGE_W),
        ]
        self.lines.append({"text": text, "box_2d": box})
        self.y += 10

    def _wrap(self, text, max_width):
        words = text.split()
        lines, current = [], ""
        for w in words:
            trial = (current + " " + w).strip()
            if self.draw.textlength(trial, font=self.font) <= max_width:
                current = trial
            else:
                if current:
                    lines.append(current)
                current = w
        if current:
            lines.append(current)
        return lines or [text]


def load_docx_form(path):
    if docx is None:
        raise RuntimeError("python-docx is not installed; cannot open Word files.")
    document = docx.Document(path)
    layout = _LineLayout(title=os.path.basename(path))
    for para in document.paragraphs:
        if para.text.strip():
            layout.add_line(para.text)
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                layout.add_line("   |   ".join(cells), indent=20)
    return layout.img, layout.lines


def load_xlsx_form(path):
    if openpyxl is None:
        raise RuntimeError("openpyxl is not installed; cannot open Excel files.")
    wb = openpyxl.load_workbook(path, data_only=True)
    sheet = wb.worksheets[0]
    layout = _LineLayout(title=f"{os.path.basename(path)} — {sheet.title}")
    max_row = min(sheet.max_row or 0, 200)
    max_col = min(sheet.max_column or 0, 20)
    for r in range(1, max_row + 1):
        cells = []
        for c in range(1, max_col + 1):
            val = sheet.cell(row=r, column=c).value
            if val not in (None, ""):
                cells.append(str(val))
        if cells:
            layout.add_line("   |   ".join(cells))
    return layout.img, layout.lines


def load_pptx_form(path, slide_number=0):
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed; cannot open PowerPoint files.")
    prs = Presentation(path)
    slide_number = max(0, min(slide_number, len(prs.slides) - 1))
    slide = prs.slides[slide_number]
    layout = _LineLayout(title=f"{os.path.basename(path)} — slide {slide_number + 1}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                layout.add_line(text)
    return layout.img, layout.lines, len(prs.slides)


# ---------------------------------------------------------------------------
# SAVED-COPY OVERLAY RENDERING — identical to the Tkinter version
# ---------------------------------------------------------------------------

def box_to_pixels(box_2d, img_w, img_h):
    ymin, xmin, ymax, xmax = box_2d
    return (
        int(xmin / 1000 * img_w), int(ymin / 1000 * img_h),
        int(xmax / 1000 * img_w), int(ymax / 1000 * img_h),
    )


def _mark_badges_only(img, fields, badge_font, language="English"):
    draw = ImageDraw.Draw(img, "RGBA")
    w, h = img.size
    for idx, field in enumerate(fields, start=1):
        x1, y1, x2, y2 = box_to_pixels(field["box_2d"], w, h)
        draw.rectangle([x1, y1, x2, y2], outline=BADGE_BG, width=3)
        badge_r = max(14, h // 80)
        bx, by = max(badge_r, x1), max(badge_r, y1 - badge_r)
        draw.ellipse([bx - badge_r, by - badge_r, bx + badge_r, by + badge_r],
                     fill=BADGE_BG, outline=(255, 255, 255, 255), width=2)
        draw.text((bx, by), local_number(idx, language), fill=BADGE_TEXT,
                   font=badge_font, anchor="mm")
    return img


def _wrap_for_note(draw, text, font, max_width):
    words = text.split()
    lines, current = [], ""
    for w in words:
        trial = (current + " " + w).strip()
        try:
            width = draw.textlength(trial, font=font)
        except Exception:
            width = len(trial) * (font.size * 0.55 if hasattr(font, "size") else 8)
        if width <= max_width or not current:
            current = trial
        else:
            lines.append(current)
            current = w
    if current:
        lines.append(current)
    return lines[:6] or [text[:40]]


def _legend_row_data(fields, start_index, col_width, note_font, header_font, tmp_draw,
                      language="English"):
    """Wrap every field's header/body/example text to fit col_width and
    measure each resulting row's height - shared by both the one-column
    and two-column legend layouts below."""
    pad = 16
    line_h = (getattr(note_font, "size", 18)) + 8
    header_line_h = (getattr(header_font, "size", 20)) + 6
    inner_w = col_width - 2 * pad

    required_word = REQUIRED_WORD.get(language, "Required")
    example_word = EXAMPLE_WORD.get(language, "Example")

    rows = []
    for offset, field in enumerate(fields):
        idx = start_index + offset
        num = local_number(idx, language)
        required = f"   {required_word}" if field.get("required") else ""
        header_text = f"{num}  {field.get('label_t', '')}{required}"
        body_text = field.get("explanation_t", "")
        example = field.get("example_t", "")
        example_text = f"{example_word}  {example}" if example else ""

        header_lines = _wrap_for_note(tmp_draw, header_text, header_font, inner_w)
        body_lines = _wrap_for_note(tmp_draw, body_text, note_font, inner_w) if body_text else []
        example_lines = _wrap_for_note(tmp_draw, example_text, note_font, inner_w) if example_text else []

        row_h = (pad + len(header_lines) * header_line_h + len(body_lines) * line_h
                 + len(example_lines) * line_h + pad)
        rows.append((header_lines, body_lines, example_lines, row_h))
    return rows, line_h, header_line_h


def _draw_legend_column(draw, rows, x_offset, col_width, line_h, header_line_h,
                         note_font, header_font):
    pad = 16
    row_gap = 10
    y = pad
    for header_lines, body_lines, example_lines, row_h in rows:
        draw.rectangle([x_offset + pad // 2, y, x_offset + col_width - pad // 2, y + row_h],
                        fill=BRIGHT_NOTE_BG, outline=BRIGHT_NOTE_BORDER, width=2)
        ty = y + 8
        tx = x_offset + pad
        for line in header_lines:
            draw.text((tx, ty), line, fill=BRIGHT_NOTE_TEXT, font=header_font)
            ty += header_line_h
        for line in body_lines:
            draw.text((tx, ty), line, fill=BRIGHT_NOTE_TEXT, font=note_font)
            ty += line_h
        for line in example_lines:
            draw.text((tx, ty), line, fill=(27, 94, 32), font=note_font)
            ty += line_h
        y += row_h + row_gap
    return y


def render_annotated_image(base_image, fields, language):
    """Builds the saved-copy image: the original form, marked with ONLY
    numbered badges (no text on top of it), placed in the CENTER of the
    canvas and flanked on the left and right by two columns of numbered
    guidance boxes - the first half of the fields on the left, the second
    half on the right. This mirrors a typical annotated-form reference
    layout (form in the middle, instructions either side) rather than
    stacking a long legend below a full-width form, which made saved
    copies of forms with many fields painfully tall to scroll through.
    """
    img = base_image.copy().convert("RGB")
    w, h = img.size
    badge_font, _ = resolve_font(language, max(16, h // 90), bold=True)
    note_font, script_ok = resolve_font(language, max(16, w // 70))
    header_font, _ = resolve_font(language, max(18, w // 60), bold=True)

    marked = _mark_badges_only(img, fields, badge_font, language=language)

    margin = 24
    col_gap = 24
    row_gap = 10
    pad = 16
    col_width = max(240, min(420, int(w * 0.34)))
    tmp = ImageDraw.Draw(Image.new("RGB", (col_width, 10), "white"))

    split = (len(fields) + 1) // 2
    left_fields, right_fields = fields[:split], fields[split:]

    if left_fields:
        left_rows, line_h, header_line_h = _legend_row_data(
            left_fields, 1, col_width, note_font, header_font, tmp, language=language)
    else:
        left_rows, line_h, header_line_h = [], 0, 0
    if right_fields:
        right_rows, r_line_h, r_header_line_h = _legend_row_data(
            right_fields, len(left_fields) + 1, col_width, note_font, header_font, tmp,
            language=language)
        if not left_fields:
            line_h, header_line_h = r_line_h, r_header_line_h
    else:
        right_rows = []

    left_h = sum(r[3] + row_gap for r in left_rows) + pad if left_rows else 0
    right_h = sum(r[3] + row_gap for r in right_rows) + pad if right_rows else 0
    has_left, has_right = bool(left_rows), bool(right_rows)

    # Scale the form down (if needed) so it reads as a centered focal
    # image between the two instruction columns instead of towering over
    # them, while never shrinking it below roughly half its original size.
    # With no columns at all, there is nothing to balance against, so the
    # form is kept at its original size.
    if has_left or has_right:
        columns_h = max(left_h, right_h, 60)
        min_form_h = int(h * 0.45)
        target_form_h = max(columns_h, min_form_h)
        scale = min(1.0, target_form_h / h)
    else:
        scale = 1.0
    scaled_w, scaled_h = max(1, int(w * scale)), max(1, int(h * scale))
    marked_scaled = marked.resize((scaled_w, scaled_h), Image.LANCZOS) if scale < 1.0 else marked

    canvas_h = max(left_h, right_h, scaled_h) + 2 * margin
    canvas_w = margin
    if has_left:
        canvas_w += col_width + col_gap
    canvas_w += scaled_w
    if has_right:
        canvas_w += col_gap + col_width
    canvas_w += margin

    canvas = Image.new("RGB", (canvas_w, canvas_h), "white")

    x = margin
    if has_left:
        left_panel = Image.new("RGB", (col_width, left_h), "white")
        _draw_legend_column(ImageDraw.Draw(left_panel), left_rows, 0, col_width,
                             line_h, header_line_h, note_font, header_font)
        canvas.paste(left_panel, (x, (canvas_h - left_h) // 2))
        x += col_width + col_gap

    canvas.paste(marked_scaled, (x, (canvas_h - scaled_h) // 2))
    x += scaled_w

    if has_right:
        x += col_gap
        right_panel = Image.new("RGB", (col_width, right_h), "white")
        _draw_legend_column(ImageDraw.Draw(right_panel), right_rows, 0, col_width,
                             line_h, header_line_h, note_font, header_font)
        canvas.paste(right_panel, (x, (canvas_h - right_h) // 2))

    return canvas, script_ok


# ---------------------------------------------------------------------------
# TEXT-TO-SPEECH — identical to the Tkinter version
# ---------------------------------------------------------------------------

def speak_text(text, language, on_error=None):
    def _run():
        try:
            if gTTS is None:
                raise RuntimeError("gTTS library is not installed.")
            lang_code = TTS_LANG_CODES.get(language, "en")
            tts = gTTS(text=text, lang=lang_code)
            fd, path = tempfile.mkstemp(suffix=".mp3")
            os.close(fd)
            tts.save(path)
            _play_audio(path)
        except Exception as exc:  # noqa: BLE001
            if on_error:
                on_error(str(exc))
    threading.Thread(target=_run, daemon=True).start()


def _play_audio(path):
    try:
        if playsound is not None:
            playsound.playsound(path)
            return
    except Exception:
        pass
    try:
        if os.name == "nt":
            os.startfile(path)  # noqa: S606
        elif sys.platform == "darwin":
            subprocess.Popen(["open", path])
        else:
            subprocess.Popen(["xdg-open", path])
    except Exception as exc:
        print("FormMitra: could not play audio:", exc)


# ---------------------------------------------------------------------------
# NATIVE FILE DIALOGS
# ---------------------------------------------------------------------------
# Kivy doesn't ship a good native "Open"/"Save As" dialog, and building a
# custom in-app file browser would be a UX downgrade, not an upgrade - so
# these use Python's built-in tkinter dialogs (a hidden root window) purely
# for the native OS file picker, while everything else in the app is Kivy.

def _native_open_dialog(title, filetypes):
    import tkinter as tk
    from tkinter import filedialog
    root = tk.Tk()
    root.withdraw()
    root.attributes("-topmost", True)
    path = filedialog.askopenfilename(title=title, filetypes=filetypes, parent=root)
    root.destroy()
    return path


def _native_save_dialog(title, default_name, filetypes, default_extension=""):
    import tkinter as tk
    from tkinter import filedialog
    root = tk.Tk()
    root.withdraw()
    root.attributes("-topmost", True)
    # defaultextension is what makes Windows/macOS append the right
    # extension when the user types a bare filename and clicks Save
    # without touching the "Save as type" dropdown - without it, a path
    # with no extension at all can come back, which later fails to save
    # with a confusing "unknown file extension:" error.
    path = filedialog.asksaveasfilename(title=title, initialfile=default_name,
                                         filetypes=filetypes, parent=root,
                                         defaultextension=default_extension)
    root.destroy()
    return path


# ---------------------------------------------------------------------------
# OVERLAY WIDGETS (numbered badges + click-to-reveal note, on the form image)
# ---------------------------------------------------------------------------

class FieldBadge(FloatLayout):
    def __init__(self, index, on_tap, active=False, language="English", font_name=None,
                 **kwargs):
        super().__init__(size_hint=(None, None), size=(dp(34), dp(34)), **kwargs)
        self.index = index
        self.on_tap = on_tap
        self.active = active
        self._ellipse = None
        with self.canvas.before:
            self._color = Color(*(BADGE_ACTIVE_RGBA if active else BADGE_RGBA))
            self._ellipse = Ellipse(pos=self.pos, size=self.size)
        # The number label's pos/size are synced explicitly on every move -
        # relying on size_hint/pos_hint alone left the label at its literal
        # default (100, 100) size in some layout contexts inside this app
        # (this FloatLayout nested inside a ScrollView, rebuilt on every
        # redraw_overlay() call), which is what made the badge numbers
        # invisible/misplaced. Explicit sync mirrors the Ellipse pattern
        # above, which has always worked reliably.
        # The badge number is shown using the same native-script digits as
        # the rest of the guidance (local_number()) - a plain ASCII "1"
        # here while every instruction box uses "૧" would be an
        # inconsistent, confusing mismatch between the number on the form
        # and the number next to its guidance. This needs the script's own
        # font (font_name), not Kivy's default, to actually draw those
        # native digit glyphs.
        self._label = MDLabel(text=local_number(index + 1, language), bold=True,
                               halign="center", valign="middle", theme_text_color="Custom",
                               text_color=(1, 1, 1, 1), font_size=sp(15),
                               size_hint=(None, None), size=self.size,
                               pos=self.pos, text_size=self.size)
        if font_name:
            self._label.font_name = font_name
        self.add_widget(self._label)
        self.bind(pos=self._sync, size=self._sync)

    def _sync(self, *_):
        self._ellipse.pos = self.pos
        self._ellipse.size = self.size
        self._label.pos = self.pos
        self._label.size = self.size
        self._label.text_size = self.size

    def on_touch_down(self, touch):
        if self.collide_point(*touch.pos):
            self.on_tap(self.index)
            return True
        return super().on_touch_down(touch)


class NoteBubble(MDLabel):
    def __init__(self, text, font_name=None, **kwargs):
        super().__init__(text=text, halign="left", valign="top",
                          theme_text_color="Custom", text_color=(0.07, 0.07, 0.07, 1),
                          bold=True, size_hint=(None, None), padding=(dp(8), dp(8)),
                          **kwargs)
        if font_name:
            self.font_name = font_name
        self._bg = None
        with self.canvas.before:
            Color(*NOTE_BG_RGBA)
            self._bg = RoundedRectangle(pos=self.pos, size=self.size, radius=[dp(6)])
            Color(*NOTE_BORDER_RGBA)
            self._line = Line(width=1.4)
        self.bind(texture_size=self._on_texture, pos=self._sync, size=self._sync)

    def _on_texture(self, *_):
        self.width = min(dp(280), max(dp(150), self.texture_size[0] + dp(16)))
        self.text_size = (self.width - dp(16), None)
        self.height = self.texture_size[1] + dp(16)

    def _sync(self, *_):
        self._bg.pos = self.pos
        self._bg.size = self.size
        self._line.rounded_rectangle = (self.x, self.y, self.width, self.height, dp(6))


# ---------------------------------------------------------------------------
# GUIDANCE CARD (side list) — a real MDCard, for a modern elevated look
# ---------------------------------------------------------------------------

class GuidanceCard(MDCard):
    def __init__(self, index, field, on_listen, on_show, active, font_name, language="English",
                 **kwargs):
        super().__init__(orientation="vertical", size_hint_y=None, padding=dp(12),
                          spacing=dp(4), radius=[dp(10)], elevation=2 if not active else 6,
                          md_bg_color=(1, 0.95, 0.46, 1) if active else (1, 1, 1, 1),
                          **kwargs)
        self.bind(minimum_height=self.setter("height"))

        req = f"   {REQUIRED_WORD.get(language, 'Required')}" if field.get("required") else ""
        num = local_number(index + 1, language)
        header = MDLabel(text=f"{num}  {field.get('label_t', '')}{req}", bold=True,
                          theme_text_color="Custom", text_color=(0.07, 0.07, 0.07, 1),
                          size_hint_y=None, halign="left", valign="top", font_style="Subtitle1")
        if font_name:
            header.font_name = font_name
        header.bind(width=lambda w, v: setattr(header, "text_size", (v, None)))
        header.bind(texture_size=lambda w, v: setattr(header, "height", v[1]))
        self.add_widget(header)

        body = MDLabel(text=field.get("explanation_t", ""), theme_text_color="Custom",
                        text_color=(0.15, 0.15, 0.15, 1), size_hint_y=None,
                        halign="left", valign="top")
        if font_name:
            body.font_name = font_name
        body.bind(width=lambda w, v: setattr(body, "text_size", (v, None)))
        body.bind(texture_size=lambda w, v: setattr(body, "height", v[1]))
        self.add_widget(body)

        example = field.get("example_t", "")
        if example:
            example_word = EXAMPLE_WORD.get(language, "Example")
            ex = MDLabel(text=f"{example_word}  {example}", theme_text_color="Custom",
                         text_color=(0.18, 0.49, 0.2, 1), italic=True, size_hint_y=None,
                         halign="left", valign="top")
            if font_name:
                ex.font_name = font_name
            ex.bind(width=lambda w, v: setattr(ex, "text_size", (v, None)))
            ex.bind(texture_size=lambda w, v: setattr(ex, "height", v[1]))
            self.add_widget(ex)

        # NOTE: MDFlatButton silently ignores its `icon` kwarg in KivyMD 1.2.0
        # (it mixes in ButtonContentsText, which has no icon widget at all) -
        # that was the cause of the missing LISTEN icon. The *IconButton
        # variants below actually render an icon glyph next to the text.
        btn_row = BoxLayout(size_hint_y=None, height=dp(44), spacing=dp(10))
        listen_btn = MDFillRoundFlatIconButton(
            text="LISTEN", icon="volume-high",
            md_bg_color=ACCENT_BLUE, text_color=(1, 1, 1, 1), icon_color=(1, 1, 1, 1),
        )
        listen_btn.bind(on_release=lambda *_: on_listen(field))
        show_btn = MDRoundFlatIconButton(
            text="SHOW ON FORM", icon="crosshairs-gps",
            line_color=ACCENT_BLUE, text_color=ACCENT_BLUE, icon_color=ACCENT_BLUE,
        )
        show_btn.bind(on_release=lambda *_: on_show(index))
        btn_row.add_widget(listen_btn)
        btn_row.add_widget(show_btn)
        btn_row.add_widget(Widget())  # spacer so the two buttons hug the left edge
        self.add_widget(btn_row)


# ---------------------------------------------------------------------------
# EMBEDDED CAMERA DIALOG (replaces the old OpenCV popup window)
# ---------------------------------------------------------------------------

class _LiveCameraWidget(KivyImage):
    """A live webcam preview built directly on OpenCV + a Kivy Texture,
    instead of Kivy's own kivy.uix.camera.Camera widget. That widget's
    OpenCV backend (CameraOpenCV) has a real bug that shows up on some
    Windows camera/driver combinations: if the very first frame read right
    after opening the device fails (which is common while the camera is
    still warming up, or briefly busy from a previous app), it crashes
    with "'CameraOpenCV' object has no attribute 'fps'" - that attribute
    is only ever set *after* the first read succeeds, and nothing catches
    the failure before then. Reading frames ourselves lets us retry
    through that warm-up window instead of crashing, and fail with a
    plain message if the camera genuinely never responds."""

    WARMUP_ATTEMPTS = 25
    WARMUP_DELAY = 0.1

    def __init__(self, index=0, resolution=(1280, 960), fps=30, **kwargs):
        super().__init__(**kwargs)
        self.last_frame = None  # most recent frame, as a right-way-up RGB PIL Image
        self._capture = None
        self._update_ev = None
        self._closed = False

        if cv2 is None:
            raise RuntimeError("OpenCV (cv2) is not installed")

        # On Windows the default backend (MSMF) is the one most prone to
        # returning a failed read for the first frame or two; DirectShow
        # is more reliable there. Elsewhere, let OpenCV pick automatically.
        if sys.platform == "win32" and hasattr(cv2, "CAP_DSHOW"):
            cap = cv2.VideoCapture(index, cv2.CAP_DSHOW)
        else:
            cap = cv2.VideoCapture(index)

        if not cap.isOpened():
            cap.release()
            raise RuntimeError("No camera found (or it's already in use by another app)")

        cap.set(cv2.CAP_PROP_FRAME_WIDTH, resolution[0])
        cap.set(cv2.CAP_PROP_FRAME_HEIGHT, resolution[1])

        # Warm-up: give the camera a few tries to deliver its first real
        # frame before giving up - this is exactly the window where Kivy's
        # own Camera widget can crash (see class docstring above).
        frame = None
        for _ in range(self.WARMUP_ATTEMPTS):
            ret, frame = cap.read()
            if ret and frame is not None:
                break
            frame = None
            time.sleep(self.WARMUP_DELAY)
        if frame is None:
            cap.release()
            raise RuntimeError("Camera did not respond - make sure no other app "
                                "(Zoom, Teams, ...) is using it, then try again")

        self._capture = cap
        self.paused = False
        self._update_ev = Clock.schedule_interval(self._update, 1.0 / max(fps, 1))

    def pause(self):
        """Stop refreshing the live texture (used while the user is
        reviewing a captured photo) without releasing the camera device,
        so resuming for a retake doesn't have to redo the warm-up wait."""
        self.paused = True

    def resume(self):
        self.paused = False

    def _update(self, dt):
        if self._closed or self._capture is None or self.paused:
            return
        ret, frame = self._capture.read()
        if not ret or frame is None:
            return
        rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
        self.last_frame = Image.fromarray(rgb)
        h, w = rgb.shape[:2]
        texture = Texture.create(size=(w, h), colorfmt="rgb")
        # Kivy textures are bottom-up; our frame array is top-down.
        texture.blit_buffer(cv2.flip(rgb, 0).tobytes(), colorfmt="rgb", bufferfmt="ubyte")
        self.texture = texture

    def stop(self):
        self._closed = True
        if self._update_ev is not None:
            self._update_ev.cancel()
            self._update_ev = None
        if self._capture is not None:
            self._capture.release()
            self._capture = None


def _pil_to_texture(pil_img):
    """Build a Kivy Texture from a right-way-up PIL RGB image, for showing
    a still (already-captured) photo in a KivyImage widget."""
    img = pil_img.convert("RGB")
    w, h = img.size
    texture = Texture.create(size=(w, h), colorfmt="rgb")
    # Kivy textures are bottom-up; PIL images are top-down.
    texture.blit_buffer(img.transpose(Image.FLIP_TOP_BOTTOM).tobytes(),
                         colorfmt="rgb", bufferfmt="ubyte")
    return texture


class CameraScreen(FloatLayout):
    """A smaller, centered camera "window" over a dimmed backdrop - rather
    than covering the whole screen, so it reads as a dialog sitting on top
    of the app instead of taking over the laptop's entire display.

    A bright guide rectangle is drawn over the live preview, sized for a
    typical A4/letter document, with the area outside it dimmed - laptop
    webcams are usually very wide-angle, so without a guide the live image
    shows a lot more desk/background than form. Aligning the document
    inside the guide box and capturing crops the photo to that box, which
    both frames the shot properly and compensates for the wide field of
    view.

    After a capture, the photo is shown full-size in a review step with
    Retake / Use This Photo buttons - rather than immediately handing an
    unclear or badly-aligned shot off to the app - so a blurry photo can
    just be retaken instead of restarting the whole scan from the toolbar."""

    # Guide box sizing, as a fraction of the camera preview's own displayed
    # image area (not the whole panel) - width fraction plus an aspect
    # ratio close to A4/letter portrait paper (~0.71 width:height).
    GUIDE_WIDTH_FRAC = 0.82
    GUIDE_ASPECT = 0.71
    GUIDE_MAX_HEIGHT_FRAC = 0.86

    TOP_BAR_HEIGHT = dp(52)

    # The camera "window" itself - a fraction of the full screen, not the
    # whole thing, per the request that opening the camera shouldn't take
    # over the entire laptop display.
    PANEL_WIDTH_FRAC = 0.6
    PANEL_HEIGHT_FRAC = 0.72

    def __init__(self, on_captured, on_cancel, **kwargs):
        super().__init__(**kwargs)
        self.on_captured = on_captured
        self.on_cancel = on_cancel
        self.camera = None
        self._guide_frac = None  # (rel_x, rel_y, rel_w, rel_h), bottom-up, 0..1
        self._mode = "live"  # or "review"
        self._pending_image = None
        self._preview_widget = None

        # Dimmed backdrop covering the whole screen, so the smaller camera
        # panel below reads as a dialog floating over the app rather than
        # a jarring full-screen takeover.
        with self.canvas.before:
            Color(0, 0, 0, 0.6)
            self._scrim = Rectangle(pos=self.pos, size=self.size)
        self.bind(pos=self._sync_scrim, size=self._sync_scrim)

        # The camera "window" itself: a smaller, centered panel.
        self.panel = FloatLayout(size_hint=(self.PANEL_WIDTH_FRAC, self.PANEL_HEIGHT_FRAC),
                                  pos_hint={"center_x": 0.5, "center_y": 0.5})
        with self.panel.canvas.before:
            Color(0.05, 0.05, 0.05, 1)
            self._panel_bg = RoundedRectangle(pos=self.panel.pos, size=self.panel.size,
                                               radius=[dp(10)])
            Color(1, 1, 1, 0.25)
            self._panel_border = Line(rounded_rectangle=(0, 0, 0, 0, dp(10)), width=dp(1.2))
        self.panel.bind(pos=self._sync_panel_bg, size=self._sync_panel_bg)
        self.add_widget(self.panel)

        try:
            self.camera = _LiveCameraWidget(resolution=(1280, 960), size_hint=(1, 1),
                                             pos_hint={"x": 0, "y": 0},
                                             allow_stretch=True, keep_ratio=True)
            self.panel.add_widget(self.camera)
            self._build_guide_graphics()
            self.camera.bind(pos=self._update_guide, size=self._update_guide,
                              texture=self._update_guide)
            Clock.schedule_once(self._update_guide, 0.3)

            self._preview_widget = KivyImage(size_hint=(1, 1), pos_hint={"x": 0, "y": 0},
                                              allow_stretch=True, keep_ratio=True, opacity=0)
            self.panel.add_widget(self._preview_widget)
        except Exception as exc:
            self.camera = None
            self.panel.add_widget(MDLabel(text=f"Could not start the camera:\n{exc}",
                                           halign="center", theme_text_color="Custom",
                                           text_color=(1, 1, 1, 1),
                                           pos_hint={"center_x": 0.5, "center_y": 0.5}))

        # -- Top bar: caption on the left, buttons on the right - anchored
        # to the top of the smaller panel (not the edge of the screen), so
        # it's always exactly where the visible "window" is.
        self._top_bar = FloatLayout(size_hint=(1, None), height=self.TOP_BAR_HEIGHT,
                                     pos_hint={"x": 0, "top": 1})
        with self._top_bar.canvas.before:
            Color(0, 0, 0, 0.55)
            self._top_bar_bg = Rectangle(pos=self._top_bar.pos, size=self._top_bar.size)
        self._top_bar.bind(pos=self._sync_top_bar_bg, size=self._sync_top_bar_bg)

        self._caption = MDLabel(
            text="", halign="left", valign="middle", theme_text_color="Custom",
            text_color=(1, 1, 1, 1), size_hint=(0.42, 1),
            pos_hint={"x": 0, "y": 0}, padding=(dp(14), 0), font_style="Caption")
        self._top_bar.add_widget(self._caption)

        btn_scroll = ScrollView(size_hint=(0.56, 1), pos_hint={"right": 1, "y": 0},
                                 do_scroll_x=True, do_scroll_y=False, bar_width=dp(3))
        self._btn_row = BoxLayout(size_hint=(None, 1), spacing=dp(8), padding=(dp(8), dp(6)))
        self._btn_row.bind(minimum_width=self._btn_row.setter("width"))
        btn_scroll.add_widget(self._btn_row)
        self._top_bar.add_widget(btn_scroll)

        self.panel.add_widget(self._top_bar)
        self._refresh_buttons()

    def _sync_scrim(self, *_):
        self._scrim.pos = self.pos
        self._scrim.size = self.size

    def _sync_panel_bg(self, *_):
        self._panel_bg.pos = self.panel.pos
        self._panel_bg.size = self.panel.size
        x, y = self.panel.pos
        w, h = self.panel.size
        self._panel_border.rounded_rectangle = (x, y, w, h, dp(10))

    def _sync_top_bar_bg(self, widget, *_):
        self._top_bar_bg.pos = widget.pos
        self._top_bar_bg.size = widget.size

    # -- Top-bar buttons: rebuilt for "live" vs "review" mode ---------------

    def _refresh_buttons(self):
        self._btn_row.clear_widgets()
        cancel_btn = MDRaisedButton(text="Cancel", md_bg_color=(0.4, 0.4, 0.4, 1))
        cancel_btn.bind(on_release=lambda *_: self.on_cancel())
        self._btn_row.add_widget(cancel_btn)

        if self._mode == "live":
            self._caption.text = "Line the document up, then tap Capture" if self.camera else ""
            if self.camera is not None:
                capture_btn = MDFillRoundFlatIconButton(
                    text="Capture", icon="camera", md_bg_color=(0.13, 0.59, 0.95, 1),
                    text_color=(1, 1, 1, 1), icon_color=(1, 1, 1, 1))
                capture_btn.bind(on_release=lambda *_: self._capture())
                self._btn_row.add_widget(capture_btn)
        else:
            self._caption.text = "Is the photo clear enough to use?"
            retake_btn = MDRaisedButton(text="Retake", md_bg_color=(0.55, 0.35, 0.05, 1))
            retake_btn.bind(on_release=lambda *_: self._retake())
            self._btn_row.add_widget(retake_btn)
            use_btn = MDFillRoundFlatIconButton(
                text="Use This Photo", icon="check", md_bg_color=(0.10, 0.6, 0.25, 1),
                text_color=(1, 1, 1, 1), icon_color=(1, 1, 1, 1))
            use_btn.bind(on_release=lambda *_: self._confirm_capture())
            self._btn_row.add_widget(use_btn)

    # -- Guide rectangle: drawn over the preview, remembered as fractions
    #    of the displayed image so capture-time cropping can reuse it -----

    def _build_guide_graphics(self):
        with self.camera.canvas.after:
            self._mask_color = Color(0, 0, 0, 0.55)
            self._mask_top = Rectangle(pos=(0, 0), size=(0, 0))
            self._mask_bottom = Rectangle(pos=(0, 0), size=(0, 0))
            self._mask_left = Rectangle(pos=(0, 0), size=(0, 0))
            self._mask_right = Rectangle(pos=(0, 0), size=(0, 0))
            Color(1, 1, 1, 0.95)
            self._guide_line = Line(rectangle=(0, 0, 0, 0), width=dp(2.5))

    def _update_guide(self, *_a):
        if self.camera is None:
            return
        dw, dh = self.camera.norm_image_size
        if dw <= 1 or dh <= 1:
            return
        img_x = self.camera.x + (self.camera.width - dw) / 2.0
        img_y = self.camera.y + (self.camera.height - dh) / 2.0

        gw = dw * self.GUIDE_WIDTH_FRAC
        gh = gw / self.GUIDE_ASPECT
        if gh > dh * self.GUIDE_MAX_HEIGHT_FRAC:
            gh = dh * self.GUIDE_MAX_HEIGHT_FRAC
            gw = gh * self.GUIDE_ASPECT
        gx = img_x + (dw - gw) / 2.0
        gy = img_y + (dh - gh) / 2.0

        # Remembered as fractions of the displayed image, y measured
        # bottom-up (Kivy's convention) - reused as-is at capture time to
        # crop the underlying frame, which shares the same aspect ratio.
        self._guide_frac = ((gx - img_x) / dw, (gy - img_y) / dh, gw / dw, gh / dh)

        self._guide_line.rectangle = (gx, gy, gw, gh)
        # Four bands covering the displayed image area outside the guide
        # box (above/below/left/right of it), dimming everything except
        # the document placement area.
        self._mask_top.pos = (img_x, gy + gh)
        self._mask_top.size = (dw, (img_y + dh) - (gy + gh))
        self._mask_bottom.pos = (img_x, img_y)
        self._mask_bottom.size = (dw, gy - img_y)
        self._mask_left.pos = (img_x, gy)
        self._mask_left.size = (gx - img_x, gh)
        self._mask_right.pos = (gx + gw, gy)
        self._mask_right.size = (img_x + dw - (gx + gw), gh)

    # -- Capture -> review (Retake / Use This Photo) -> confirm -------------

    def _capture(self):
        if self.camera is None or self.camera.last_frame is None:
            return
        try:
            cropped = self._crop_to_guide(self.camera.last_frame)
        except Exception as exc:
            print("FormMitra: capture failed:", exc)
            return
        self._pending_image = cropped
        self.camera.pause()
        self._preview_widget.texture = _pil_to_texture(cropped)
        self._preview_widget.opacity = 1
        self._mode = "review"
        self._refresh_buttons()

    def _retake(self):
        self._pending_image = None
        self._preview_widget.opacity = 0
        if self.camera is not None:
            self.camera.resume()
        self._mode = "live"
        self._refresh_buttons()

    def _confirm_capture(self):
        img = self._pending_image
        self._pending_image = None
        if img is not None:
            self.on_captured(img)

    def _crop_to_guide(self, frame):
        """Crop the captured frame to the guide box the user aligned the
        document in - this is what compensates for a very wide webcam
        field of view, since only the guided region is kept."""
        if not self._guide_frac:
            return frame.copy()
        rel_x, rel_y, rel_w, rel_h = self._guide_frac
        fw, fh = frame.size
        left = int(rel_x * fw)
        right = int((rel_x + rel_w) * fw)
        # Kivy's y-axis is bottom-up; PIL's is top-down.
        top = int((1 - rel_y - rel_h) * fh)
        bottom = int((1 - rel_y) * fh)
        left, top = max(0, left), max(0, top)
        right, bottom = min(fw, right), min(fh, bottom)
        if right - left < 10 or bottom - top < 10:
            return frame.copy()
        return frame.crop((left, top, right, bottom))


# ---------------------------------------------------------------------------
# MAIN APP
# ---------------------------------------------------------------------------

class FormMitraApp(MDApp):
    def build(self):
        self.title = f"{APP_NAME} — Understand Any Form, In Your Language"
        self.theme_cls.primary_palette = "Teal"
        self.theme_cls.theme_style = "Light"

        self.cfg = load_config()
        self.api_key = self.cfg.get("api_key", "")
        self.language = self.cfg.get("last_language", "Hindi")
        self.font_scale = 1.0
        self.high_contrast = False
        self.show_all_notes = False

        self.base_image = None
        self.canonical_fields = []
        self.display_fields = []
        self.current_source_path = None
        self.active_field = None
        self.field_boxes = []
        self.overlay_widgets = []

        self.job_queue = queue.Queue()
        self._pending_action = None
        self._camera_dialog = None
        self._language_menu = None

        root = MDBoxLayout(orientation="vertical")

        # -- Top app bar --------------------------------------------------
        self.top_bar = MDTopAppBar(
            title=APP_NAME,
            right_action_items=[
                ["cog", lambda x: self.open_settings()],
                ["theme-light-dark", lambda x: self.toggle_theme()],
            ],
        )
        root.add_widget(self.top_bar)

        # -- Toolbar ----------------------------------------------------------
        # Only the two most-used actions (Upload, Camera) plus the language
        # picker are here - everything else lives in the options ribbon
        # (below) so this row never has more than 3 buttons in it and can't
        # overflow/overlap regardless of window width.
        toolbar = MDBoxLayout(size_hint_y=None, height=dp(52), padding=(dp(8), dp(6)),
                               spacing=dp(8), md_bg_color=SURFACE_TINT)

        upload_btn = MDRaisedButton(text="Upload", icon="file-upload",
                                     md_bg_color=ACCENT_TEAL)
        upload_btn.bind(on_release=lambda *_: self.on_upload())
        camera_btn = MDRaisedButton(text="Camera", icon="camera",
                                     md_bg_color=ACCENT_TEAL)
        camera_btn.bind(on_release=lambda *_: self.on_camera())
        toolbar.add_widget(upload_btn)
        toolbar.add_widget(camera_btn)

        self.lang_btn = MDRaisedButton(text=self.language, icon="translate",
                                        md_bg_color=ACCENT_TEAL_DARK)
        self.lang_btn.bind(on_release=lambda *_: self._open_language_menu())
        toolbar.add_widget(self.lang_btn)
        toolbar.add_widget(Widget())  # spacer: keeps the row left-aligned
        root.add_widget(toolbar)

        # -- Options ribbon -----------------------------------------------
        # Every secondary action (save/download, read aloud, font size,
        # high contrast, show-all-notes, theme) as its own always-visible
        # button, per the user's request to keep these visible as a ribbon
        # instead of tucked behind a "..." menu. Wrapped in a horizontal
        # ScrollView so that on a narrow window the row scrolls sideways
        # instead of ever overlapping itself - the bug that prompted the
        # "..." menu in the first place.
        ribbon_scroll = ScrollView(size_hint_y=None, height=dp(48),
                                    do_scroll_x=True, do_scroll_y=False,
                                    bar_width=dp(4))
        self.options_ribbon = MDBoxLayout(size_hint_x=None, size_hint_y=1,
                                           padding=(dp(8), dp(4)), spacing=dp(8),
                                           md_bg_color=SURFACE_TINT)
        self.options_ribbon.bind(minimum_width=self.options_ribbon.setter("width"))

        # NOTE: MDRaisedButton (like plain MDFlatButton) has no icon slot in
        # KivyMD 1.2.0 - an icon= argument is silently accepted and never
        # drawn (the same class of bug fixed for the LISTEN button earlier).
        # Every ribbon button that needs a visible icon uses the icon-
        # capable MDFillRoundFlatIconButton instead, colored to match.
        def _ribbon_btn(text, icon):
            return MDFillRoundFlatIconButton(
                text=text, icon=icon, md_bg_color=ACCENT_TEAL_DARK,
                text_color=(1, 1, 1, 1), icon_color=(1, 1, 1, 1))

        save_ribbon_btn = _ribbon_btn("Save / Download", "download")
        save_ribbon_btn.bind(on_release=lambda *_: self.on_save_copy())
        read_all_ribbon_btn = _ribbon_btn("Read All Aloud", "volume-high")
        read_all_ribbon_btn.bind(on_release=lambda *_: self.on_read_all())
        dec_font_btn = MDRaisedButton(text="A-", md_bg_color=ACCENT_TEAL_DARK)
        dec_font_btn.bind(on_release=lambda *_: self.adjust_font(-1))
        inc_font_btn = MDRaisedButton(text="A+", md_bg_color=ACCENT_TEAL_DARK)
        inc_font_btn.bind(on_release=lambda *_: self.adjust_font(1))
        self.contrast_btn = _ribbon_btn("High Contrast: Off", "checkbox-blank-outline")
        self.contrast_btn.bind(on_release=lambda *_: self.set_high_contrast(not self.high_contrast))
        self.show_all_btn = _ribbon_btn("Show All Notes: Off", "checkbox-blank-outline")
        self.show_all_btn.bind(on_release=lambda *_: self.set_show_all(not self.show_all_notes))
        theme_btn = _ribbon_btn("Theme", "theme-light-dark")
        theme_btn.bind(on_release=lambda *_: self.toggle_theme())

        for w in (save_ribbon_btn, read_all_ribbon_btn, dec_font_btn, inc_font_btn,
                  self.contrast_btn, self.show_all_btn, theme_btn):
            self.options_ribbon.add_widget(w)
        ribbon_scroll.add_widget(self.options_ribbon)
        root.add_widget(ribbon_scroll)

        hint = MDLabel(
            text="Tip: tap any numbered badge on the form (or \"Show on form\" on a card) "
                 "to reveal just that field's note. Tap it again to hide it.",
            size_hint_y=None, height=dp(24), theme_text_color="Secondary",
            font_style="Caption",
        )
        root.add_widget(hint)

        self.status_label = MDLabel(text="Upload a form or use the camera to begin.",
                                     size_hint_y=None, height=dp(28), theme_text_color="Secondary")
        root.add_widget(self.status_label)

        # -- Body: form view (left) + guidance cards (right) ------------------
        body = MDBoxLayout(orientation="horizontal", spacing=dp(4))

        left = MDBoxLayout(orientation="vertical", size_hint_x=0.62)
        self.form_scroll = ScrollView()
        self.form_area = FloatLayout(size_hint=(None, None))
        self.form_image = KivyImage(allow_stretch=True, keep_ratio=True, size_hint=(1, 1))
        self.form_area.add_widget(self.form_image)
        self.form_scroll.add_widget(self.form_area)
        left.add_widget(self.form_scroll)

        save_btn = MDRaisedButton(text="Save Annotated Copy", icon="content-save",
                                   size_hint_y=None, height=dp(44))
        save_btn.bind(on_release=lambda *_: self.on_save_copy())
        read_all_btn = MDRaisedButton(text="Read All Aloud", icon="volume-high",
                                       size_hint_y=None, height=dp(44))
        read_all_btn.bind(on_release=lambda *_: self.on_read_all())
        bottom_row = BoxLayout(size_hint_y=None, height=dp(52), spacing=dp(8), padding=dp(4))
        bottom_row.add_widget(save_btn)
        bottom_row.add_widget(read_all_btn)
        left.add_widget(bottom_row)
        body.add_widget(left)

        right = MDBoxLayout(orientation="vertical", size_hint_x=0.38)
        self.cards_scroll = ScrollView()
        self.cards_grid = GridLayout(cols=1, size_hint_y=None, spacing=dp(8), padding=dp(8))
        self.cards_grid.bind(minimum_height=self.cards_grid.setter("height"))
        self.cards_scroll.add_widget(self.cards_grid)
        right.add_widget(self.cards_scroll)
        body.add_widget(right)

        root.add_widget(body)

        Clock.schedule_once(lambda dt: self._poll_jobs(), 0.2)
        prefetch_font(self.language)
        if not self.api_key:
            Clock.schedule_once(lambda dt: self.open_settings(), 0.4)

        # The camera dialog needs to be a TRUE full-window overlay so its
        # Cancel/Capture buttons are always reachable - adding it directly
        # to `root` (a vertical MDBoxLayout) would make it just another row
        # sharing height with the form/cards area below, squeezing it down
        # to a sliver where the buttons could end up unreachable. Wrapping
        # everything in an outer FloatLayout and adding the camera dialog
        # there instead (see on_camera()) makes it fill the whole window on
        # top of everything else, the way a full-screen dialog should.
        overlay_root = FloatLayout()
        overlay_root.add_widget(root)
        return overlay_root

    # -- Theme / appearance -------------------------------------------------

    def toggle_theme(self):
        self.theme_cls.theme_style = "Dark" if self.theme_cls.theme_style == "Light" else "Light"

    def set_high_contrast(self, value):
        self.high_contrast = value
        self.theme_cls.theme_style = "Dark" if value else "Light"
        self.contrast_btn.text = f"High Contrast: {'On' if value else 'Off'}"
        self.contrast_btn.icon = "check" if value else "checkbox-blank-outline"
        self.rebuild_cards()

    def set_show_all(self, value):
        self.show_all_notes = value
        self.active_field = None
        self.show_all_btn.text = f"Show All Notes: {'On' if value else 'Off'}"
        self.show_all_btn.icon = "check" if value else "checkbox-blank-outline"
        self.redraw_overlay()

    def adjust_font(self, delta):
        self.font_scale = max(0.7, min(2.0, self.font_scale + delta * 0.15))
        self.redraw_overlay()
        self.rebuild_cards()

    # -- Language picker menu ------------------------------------------------

    def _open_language_menu(self):
        items = [
            {"text": lang, "on_release": (lambda l=lang: self._select_language(l))}
            for lang in LANGUAGES
        ]
        menu = MDDropdownMenu(caller=self.lang_btn, items=items, width_mult=4,
                               max_height=dp(400))
        self._language_menu = menu
        menu.open()

    def _select_language(self, lang):
        if self._language_menu:
            self._language_menu.dismiss()
        self.lang_btn.text = lang
        self._apply_language_change(lang)

    # -- Settings / API key --------------------------------------------------

    def open_settings(self):
        self.api_key_field = MDTextField(text=self.api_key, hint_text="Gemini API key",
                                          size_hint_x=1)
        content = MDBoxLayout(orientation="vertical", spacing=dp(12), size_hint_y=None,
                               height=dp(140), padding=dp(8))
        content.add_widget(MDLabel(
            text="Enter your free Google AI Studio (Gemini) API key.\n"
                 "Get one at aistudio.google.com/app/apikey\nSaved only on this computer.",
            size_hint_y=None, height=dp(70)))
        content.add_widget(self.api_key_field)

        self.settings_dialog = MDDialog(
            title="Settings",
            type="custom",
            content_cls=content,
            buttons=[
                MDFlatButton(text="CANCEL", on_release=lambda *_: self.settings_dialog.dismiss()),
                MDFlatButton(text="SAVE", on_release=lambda *_: self._save_api_key()),
            ],
        )
        self.settings_dialog.open()

    def _save_api_key(self):
        self.api_key = self.api_key_field.text.strip()
        self.cfg["api_key"] = self.api_key
        save_config(self.cfg)
        self.settings_dialog.dismiss()

    def _require_client(self):
        if not self.api_key:
            raise GeminiError("No Gemini API key was provided. Open Settings and paste one in.")
        return GeminiClient(self.api_key)

    def _toast(self, message):
        Snackbar(text=message).open()

    def _error_dialog(self, message):
        dialog = MDDialog(title=APP_NAME, text=message,
                           buttons=[MDFlatButton(text="OK", on_release=lambda *_: dialog.dismiss())])
        dialog.open()

    # -- Background job plumbing (mirrors the Tkinter version's queue) -----

    def _run_async(self, label, action, fn, *args, **kwargs):
        self.status_label.text = label
        self._pending_action = action

        def worker():
            try:
                result = fn(*args, **kwargs)
                self.job_queue.put(("ok", result))
            except Exception as exc:  # noqa: BLE001
                traceback.print_exc()
                self.job_queue.put(("error", str(exc)))

        threading.Thread(target=worker, daemon=True).start()

    def _poll_jobs(self):
        try:
            while True:
                kind, payload = self.job_queue.get_nowait()
                self._handle_job_result(kind, payload)
        except queue.Empty:
            pass
        Clock.schedule_once(lambda dt: self._poll_jobs(), 0.2)

    def _handle_job_result(self, kind, payload):
        action = self._pending_action
        self._pending_action = None
        if kind == "error":
            self.status_label.text = "Something went wrong."
            self._error_dialog(payload)
            return
        if action == "load_and_detect":
            self._on_form_loaded(payload)
        elif action == "translate":
            self._on_translated(payload)

    # -- Upload / camera --------------------------------------------------

    def on_upload(self):
        path = _native_open_dialog(
            "Choose a blank form",
            [("Supported forms", "*.pdf *.png *.jpg *.jpeg *.bmp *.tiff *.docx *.xlsx *.pptx"),
             ("PDF", "*.pdf"), ("Images", "*.png *.jpg *.jpeg *.bmp *.tiff"),
             ("Word", "*.docx"), ("Excel", "*.xlsx"), ("PowerPoint", "*.pptx"),
             ("All files", "*.*")],
        )
        if not path:
            return
        self.current_source_path = path
        self._run_async(f"Reading {os.path.basename(path)} ...", "load_and_detect",
                         self._load_and_detect, path)

    def on_camera(self):
        if self._camera_dialog is not None:
            return
        self._camera_dialog = CameraScreen(on_captured=self._on_captured,
                                            on_cancel=self._close_camera)
        self.root.add_widget(self._camera_dialog)

    def _close_camera(self):
        if self._camera_dialog is not None:
            if self._camera_dialog.camera is not None:
                self._camera_dialog.camera.stop()
            self.root.remove_widget(self._camera_dialog)
            self._camera_dialog = None

    def _on_captured(self, pil_image):
        self._close_camera()
        self.current_source_path = None
        self._run_async("Reading the captured photo ...", "load_and_detect",
                         self._load_and_detect_image, pil_image)

    # -- Loading + detection (runs on a worker thread) -----------------------

    def _load_and_detect(self, path):
        ext = os.path.splitext(path)[1].lower()
        client = self._require_client()

        if ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff"):
            img = load_image_form(path)
            fields = client.detect_fields_from_image(img)
        elif ext == ".pdf":
            img, _ = load_pdf_form(path, 0)
            fields = client.detect_fields_from_image(img)
        elif ext == ".docx":
            img, lines = load_docx_form(path)
            fields = client.detect_fields_from_lines(lines)
        elif ext == ".xlsx":
            img, lines = load_xlsx_form(path)
            fields = client.detect_fields_from_lines(lines)
        elif ext == ".pptx":
            img, lines, _ = load_pptx_form(path, 0)
            fields = client.detect_fields_from_lines(lines)
        else:
            raise RuntimeError(f"Unsupported file type: {ext}")

        translated = client.translate_fields(fields, self.language)
        return img, fields, translated

    def _load_and_detect_image(self, pil_image):
        client = self._require_client()
        fields = client.detect_fields_from_image(pil_image)
        translated = client.translate_fields(fields, self.language)
        return pil_image, fields, translated

    def _on_form_loaded(self, payload):
        img, canonical, translated = payload
        self.base_image = img
        self.canonical_fields = canonical
        self.display_fields = translated
        self.active_field = None
        self.status_label.text = f"Found {len(canonical)} field(s). Showing guidance in {self.language}."
        self._load_image_into_view()
        self.redraw_overlay()
        self.rebuild_cards()

    def _load_image_into_view(self):
        fd, path = tempfile.mkstemp(suffix=".png")
        os.close(fd)
        self.base_image.save(path)
        self._current_preview_path = path
        self.form_image.source = path
        self.form_image.reload()

    # -- Language change ------------------------------------------------

    def _apply_language_change(self, text):
        if text == self.language:
            return
        self.language = text
        self.cfg["last_language"] = text
        save_config(self.cfg)
        self.active_field = None
        prefetch_font(text)
        if not self.canonical_fields:
            return
        self._run_async(f"Translating guidance into {text} ...", "translate",
                         self._translate_only)

    def _translate_only(self):
        client = self._require_client()
        return client.translate_fields(self.canonical_fields, self.language)

    def _on_translated(self, translated):
        self.display_fields = translated
        self.status_label.text = f"Guidance shown in {self.language}."
        self.redraw_overlay()
        self.rebuild_cards()

    # -- Overlay: numbered badges, click-to-reveal note ----------------------

    def redraw_overlay(self):
        for w in list(self.overlay_widgets):
            self.form_area.remove_widget(w)
        self.overlay_widgets = []
        self.field_boxes = []

        if self.base_image is None:
            return

        avail_w = self.form_scroll.width or dp(500)
        img_w, img_h = self.base_image.size
        scale = avail_w / max(1, img_w)
        self.form_area.size = (avail_w, img_h * scale)
        self.form_image.size = self.form_area.size

        font_name = kivy_font_name(self.language, bold=True)
        show_all = self.show_all_notes

        for i, field in enumerate(self.display_fields):
            ymin, xmin, ymax, xmax = field["box_2d"]
            self.field_boxes.append((xmin, ymin, xmax, ymax))
            active = show_all or (self.active_field == i)

            fx = xmin / 1000.0
            fy = 1.0 - (ymax / 1000.0)
            badge = FieldBadge(i, self._toggle_field, active=active,
                                language=self.language, font_name=font_name,
                                pos_hint={"x": fx, "y": fy})
            self.form_area.add_widget(badge)
            self.overlay_widgets.append(badge)

            if active:
                note_text = field.get("explanation_t") or field.get("label_t") or ""
                num = local_number(i + 1, self.language)
                bubble = NoteBubble(f"{num}  {note_text}", font_name=font_name,
                                     pos_hint={"x": min(0.95, fx + 0.05), "y": fy})
                self.form_area.add_widget(bubble)
                self.overlay_widgets.append(bubble)

    def _toggle_field(self, i):
        self.active_field = None if self.active_field == i else i
        self.redraw_overlay()
        self.rebuild_cards()

    def _show_field(self, i):
        self.active_field = i
        self.redraw_overlay()
        self.rebuild_cards()

    # -- Guidance cards -----------------------------------------------------

    def rebuild_cards(self):
        self.cards_grid.clear_widgets()
        font_name = kivy_font_name(self.language)
        if not self.display_fields:
            self.cards_grid.add_widget(MDLabel(text="No form loaded yet.", size_hint_y=None,
                                                height=dp(40)))
            return
        for i, field in enumerate(self.display_fields):
            card = GuidanceCard(i, field, on_listen=self._listen, on_show=self._show_field,
                                 active=(self.active_field == i), font_name=font_name,
                                 language=self.language)
            self.cards_grid.add_widget(card)

    # -- Text-to-speech -------------------------------------------------------

    def _listen(self, field):
        text = f"{field.get('label_t','')}. {field.get('explanation_t','')}"
        self.status_label.text = "Speaking ..."
        speak_text(text, self.language,
                   on_error=lambda msg: self._toast(f"Could not play audio: {msg}"))

    def on_read_all(self):
        if not self.display_fields:
            self._toast("Load a form first.")
            return
        combined = ". ".join(
            f"{i}. {f.get('label_t','')}. {f.get('explanation_t','')}"
            for i, f in enumerate(self.display_fields, start=1)
        )
        self.status_label.text = "Reading all guidance aloud ..."
        speak_text(combined, self.language,
                   on_error=lambda msg: self._toast(f"Could not play audio: {msg}"))

    # -- Save annotated copy --------------------------------------------------

    def on_save_copy(self):
        if self.base_image is None or not self.display_fields:
            self._toast("Load a form and wait for guidance first.")
            return
        default_name = "form_with_guidance.png"
        if self.current_source_path:
            stem = os.path.splitext(os.path.basename(self.current_source_path))[0]
            default_name = f"{stem}_guidance.png"
        path = _native_save_dialog(
            "Save annotated copy", default_name,
            [("PNG image", "*.png"), ("PDF document", "*.pdf")],
            default_extension=".png",
        )
        if not path:
            return
        # Belt-and-braces on top of defaultextension above: some dialog
        # backends (and users who type over the whole filename) can still
        # hand back a path with no extension, or one that's neither of the
        # two we offer - default it to .png rather than letting Pillow's
        # save() raise a confusing "unknown file extension:" error.
        ext = os.path.splitext(path)[1].lower()
        if ext not in (".png", ".pdf"):
            path += ".png"
        try:
            annotated, script_ok = render_annotated_image(
                self.base_image, self.display_fields, self.language
            )
            if path.lower().endswith(".pdf"):
                annotated.save(path, "PDF", resolution=150.0)
            else:
                annotated.save(path)
            msg = f"Saved to:\n{path}"
            if not script_ok and self.language != "English":
                msg += ("\n\nNote: this computer may not have a font installed for "
                        f"{self.language}, so the text baked into the saved image "
                        "might not display correctly. The on-screen view and side list "
                        "are unaffected.")
            self._error_dialog(msg)
        except Exception as exc:
            self._error_dialog(f"Could not save file: {exc}")


# ---------------------------------------------------------------------------
# ENTRY POINT
# ---------------------------------------------------------------------------

def main():
    if Image is None:
        print("FATAL: Pillow could not be installed/imported. Please install it manually:")
        print(f"   {sys.executable} -m pip install Pillow")
        sys.exit(1)
    FormMitraApp().run()


if __name__ == "__main__":
    main()

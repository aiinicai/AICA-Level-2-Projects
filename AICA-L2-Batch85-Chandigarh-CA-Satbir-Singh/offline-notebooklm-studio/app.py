"""
Offline NotebookLM Studio - FastAPI Backend
Connects to local LM Studio (http://localhost:1234/v1) for grounded document QA,
rich document generation (DOCX, PPTX), and 3-panel source-grounded research.
"""

import os
import io
import json
import uuid
import time
import shutil
import csv
from datetime import datetime
from typing import List, Optional, Dict, Any

import uvicorn
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Query, Response
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, FileResponse, HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
import requests

# Document extractors
import pypdf
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

import pptx
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN

import openpyxl
from PIL import Image

# OpenAI client for LM Studio
from openai import OpenAI

# Initialize FastAPI App
app = FastAPI(
    title="Offline NotebookLM Studio",
    description="A 3-panel NotebookLM-style desktop application connected to local LM Studio",
    version="1.0.0"
)

# Enable CORS for local desktop usage and web clients
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Directories
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
SOURCES_DIR = os.path.join(BASE_DIR, "sources")
FILES_DIR = os.path.join(SOURCES_DIR, "files")
METADATA_FILE = os.path.join(SOURCES_DIR, "sources.json")
EXPORTS_DIR = os.path.join(BASE_DIR, "exports")

os.makedirs(FILES_DIR, exist_ok=True)
os.makedirs(EXPORTS_DIR, exist_ok=True)

# Default LM Studio Base URL
DEFAULT_LM_STUDIO_URL = os.getenv("LM_STUDIO_URL", "http://localhost:1234/v1")
DEFAULT_API_KEY = "lm-studio"


# ==========================================
# METADATA HELPERS
# ==========================================

def load_sources_metadata() -> List[Dict[str, Any]]:
    if not os.path.exists(METADATA_FILE):
        return []
    try:
        with open(METADATA_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return []

def save_sources_metadata(sources: List[Dict[str, Any]]):
    with open(METADATA_FILE, "w", encoding="utf-8") as f:
        json.dump(sources, f, indent=2, ensure_ascii=False)

def get_source_text_path(source_id: str) -> str:
    return os.path.join(FILES_DIR, f"{source_id}.txt")

def get_source_original_path(source_id: str, original_filename: str) -> str:
    ext = os.path.splitext(original_filename)[1]
    return os.path.join(FILES_DIR, f"{source_id}{ext}")


# ==========================================
# FILE TEXT EXTRACTORS
# ==========================================

def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    text = ""

    if ext == ".pdf":
        try:
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            pages_text = []
            for idx, page in enumerate(reader.pages):
                page_content = page.extract_text() or ""
                if page_content.strip():
                    pages_text.append(f"--- Page {idx + 1} ---\n{page_content.strip()}")
            text = "\n\n".join(pages_text)
            if not text.strip():
                text = f"[PDF Document '{filename}' with {len(reader.pages)} pages: No selectable text layer found.]"
        except Exception as e:
            text = f"[Error extracting PDF text from '{filename}': {str(e)}]"

    elif ext in [".docx", ".doc"]:
        try:
            doc = Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            
            # Extract tables if any
            table_texts = []
            for t_idx, table in enumerate(doc.tables):
                t_rows = []
                for row in table.rows:
                    t_rows.append(" | ".join([c.text.strip() for c in row.cells]))
                if t_rows:
                    table_texts.append(f"Table {t_idx + 1}:\n" + "\n".join(t_rows))
            
            all_parts = paragraphs + table_texts
            text = "\n\n".join(all_parts)
        except Exception as e:
            text = f"[Error extracting Word doc text from '{filename}': {str(e)}]"

    elif ext in [".xlsx", ".xls"]:
        try:
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            sheets_text = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_vals = [str(val) if val is not None else "" for val in row]
                    if any(v.strip() for v in row_vals):
                        rows.append(" | ".join(row_vals))
                if rows:
                    sheets_text.append(f"### Sheet: {sheet_name}\n" + "\n".join(rows))
            text = "\n\n".join(sheets_text)
        except Exception as e:
            text = f"[Error extracting Excel text from '{filename}': {str(e)}]"

    elif ext == ".csv":
        try:
            decoded = None
            for enc in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
                try:
                    decoded = file_bytes.decode(enc)
                    break
                except UnicodeDecodeError:
                    continue
            if decoded is None:
                decoded = file_bytes.decode("utf-8", errors="replace")
            
            reader = csv.reader(io.StringIO(decoded))
            rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
            text = "\n".join(rows)
        except Exception as e:
            text = f"[Error extracting CSV from '{filename}': {str(e)}]"

    elif ext in [".txt", ".md", ".json", ".py", ".js", ".html", ".css", ".xml", ".yaml", ".yml", ".tsv"]:
        try:
            for enc in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
                try:
                    text = file_bytes.decode(enc)
                    break
                except UnicodeDecodeError:
                    continue
            if not text:
                text = file_bytes.decode("utf-8", errors="replace")
        except Exception as e:
            text = f"[Error decoding text file '{filename}': {str(e)}]"

    elif ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif"]:
        try:
            img = Image.open(io.BytesIO(file_bytes))
            w, h = img.size
            format_name = img.format or ext.replace(".", "").upper()
            mode = img.mode
            
            exif_info = ""
            if hasattr(img, "_getexif") and img._getexif():
                exif_info = f" (EXIF metadata present, {len(img._getexif())} tags)"
                
            text = f"[Image Source: {filename}]\n- Dimensions: {w} x {h} px\n- Format: {format_name}\n- Color Mode: {mode}{exif_info}\n- Note: Use this source as visual reference context."
        except Exception as e:
            text = f"[Image info for '{filename}': {str(e)}]"

    else:
        # Fallback raw text attempt
        try:
            text = file_bytes.decode("utf-8", errors="ignore")
        except Exception:
            text = f"[Binary file '{filename}' of format {ext}]"

    return text.strip()


# ==========================================
# PYDANTIC SCHEMAS
# ==========================================

class RenameSourceRequest(BaseModel):
    new_name: str

class ChatMessage(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    messages: List[ChatMessage]
    selected_source_ids: List[str] = []
    temperature: float = 0.7
    model: Optional[str] = None
    system_prompt: Optional[str] = None
    lm_studio_url: Optional[str] = None

class ExportDocxRequest(BaseModel):
    markdown_content: str
    filename: Optional[str] = "NotebookLM_Synthesis.docx"
    title: Optional[str] = "Document Synthesis & Analysis"

class ExportPptxRequest(BaseModel):
    markdown_content: str
    filename: Optional[str] = "NotebookLM_Presentation.pptx"
    title: Optional[str] = "Executive Presentation"

class UseOutputAsSourceRequest(BaseModel):
    content: str
    title: Optional[str] = "Output Synthesis"


# ==========================================
# LM STUDIO HEALTH & MODELS
# ==========================================

@app.get("/api/lm-status")
def check_lm_studio_status(url: Optional[str] = None):
    base_url = (url or DEFAULT_LM_STUDIO_URL).rstrip("/")
    if not base_url.endswith("/v1"):
        models_url = f"{base_url}/v1/models"
    else:
        models_url = f"{base_url}/models"

    try:
        resp = requests.get(models_url, timeout=3)
        if resp.status_code == 200:
            data = resp.json()
            models = [m.get("id") for m in data.get("data", [])] if "data" in data else []
            return {
                "status": "connected",
                "base_url": base_url,
                "models": models,
                "active_model": models[0] if models else "default-model",
                "message": f"Connected to LM Studio ({len(models)} model(s) available)"
            }
        else:
            return {
                "status": "error",
                "base_url": base_url,
                "status_code": resp.status_code,
                "models": [],
                "message": f"LM Studio responded with HTTP {resp.status_code}"
            }
    except requests.exceptions.ConnectionError:
        return {
            "status": "disconnected",
            "base_url": base_url,
            "models": [],
            "message": "Cannot connect to LM Studio. Ensure LM Studio is open with Local Server running on port 1234."
        }
    except Exception as e:
        return {
            "status": "error",
            "base_url": base_url,
            "models": [],
            "message": str(e)
        }


# ==========================================
# SOURCE MANAGEMENT ENDPOINTS
# ==========================================

@app.get("/api/sources")
def list_sources():
    sources = load_sources_metadata()
    # Sort by created_at descending
    sources.sort(key=lambda s: s.get("created_at", ""), reverse=True)
    return {"sources": sources}

@app.post("/api/upload")
async def upload_sources(files: List[UploadFile] = File(...)):
    existing_sources = load_sources_metadata()
    newly_created = []

    for file in files:
        file_bytes = await file.read()
        filename = file.filename or "untitled_source"
        source_id = str(uuid.uuid4())[:12]
        ext = os.path.splitext(filename)[1].lower().replace(".", "") or "txt"
        
        # Save raw original file
        orig_path = get_source_original_path(source_id, filename)
        with open(orig_path, "wb") as f:
            f.write(file_bytes)

        # Extract text
        extracted_text = extract_text_from_file(file_bytes, filename)
        
        # Save extracted text
        text_path = get_source_text_path(source_id)
        with open(text_path, "w", encoding="utf-8") as f:
            f.write(extracted_text)

        # Preview text snippet (first 300 chars)
        preview = extracted_text[:300].strip().replace("\n", " ")
        if len(extracted_text) > 300:
            preview += "..."

        source_info = {
            "id": source_id,
            "name": filename,
            "file_type": ext.upper(),
            "size_bytes": len(file_bytes),
            "char_count": len(extracted_text),
            "preview": preview,
            "created_at": datetime.now().isoformat()
        }
        
        existing_sources.append(source_info)
        newly_created.append(source_info)

    save_sources_metadata(existing_sources)
    return {"message": f"Uploaded {len(newly_created)} source(s)", "sources": newly_created}

@app.get("/api/sources/{source_id}/text")
def get_source_text(source_id: str):
    sources = load_sources_metadata()
    source = next((s for s in sources if s["id"] == source_id), None)
    if not source:
        raise HTTPException(status_code=404, detail="Source not found")

    text_path = get_source_text_path(source_id)
    if not os.path.exists(text_path):
        return {"id": source_id, "name": source["name"], "text": "[No text extracted]"}

    with open(text_path, "r", encoding="utf-8") as f:
        full_text = f.read()

    return {
        "id": source_id,
        "name": source["name"],
        "file_type": source.get("file_type", "TXT"),
        "char_count": len(full_text),
        "text": full_text
    }

@app.put("/api/sources/{source_id}/rename")
def rename_source(source_id: str, req: RenameSourceRequest):
    sources = load_sources_metadata()
    source = next((s for s in sources if s["id"] == source_id), None)
    if not source:
        raise HTTPException(status_code=404, detail="Source not found")

    old_name = source["name"]
    source["name"] = req.new_name.strip() or old_name
    save_sources_metadata(sources)
    return {"message": "Source renamed successfully", "source": source}

@app.delete("/api/sources/{source_id}")
def delete_source(source_id: str):
    sources = load_sources_metadata()
    source = next((s for s in sources if s["id"] == source_id), None)
    if not source:
        raise HTTPException(status_code=404, detail="Source not found")

    # Remove files
    text_path = get_source_text_path(source_id)
    if os.path.exists(text_path):
        os.remove(text_path)
    
    orig_path = get_source_original_path(source_id, source["name"])
    if os.path.exists(orig_path):
        os.remove(orig_path)

    updated_sources = [s for s in sources if s["id"] != source_id]
    save_sources_metadata(updated_sources)
    return {"message": "Source deleted successfully", "deleted_id": source_id}

@app.post("/api/sources/from-output")
def add_output_as_source(req: UseOutputAsSourceRequest):
    existing_sources = load_sources_metadata()
    source_id = str(uuid.uuid4())[:12]
    
    clean_title = req.title.strip() if req.title else "Generated Synthesis"
    if not clean_title.endswith(".md"):
        filename = f"{clean_title}.md"
    else:
        filename = clean_title
        
    text_content = req.content.strip()
    
    # Save text file
    text_path = get_source_text_path(source_id)
    with open(text_path, "w", encoding="utf-8") as f:
        f.write(text_content)

    # Save original md
    orig_path = get_source_original_path(source_id, filename)
    with open(orig_path, "w", encoding="utf-8") as f:
        f.write(text_content)

    preview = text_content[:300].strip().replace("\n", " ")
    if len(text_content) > 300:
        preview += "..."

    new_source = {
        "id": source_id,
        "name": filename,
        "file_type": "MD",
        "size_bytes": len(text_content.encode("utf-8")),
        "char_count": len(text_content),
        "preview": preview,
        "created_at": datetime.now().isoformat(),
        "is_generated": True
    }

    existing_sources.append(new_source)
    save_sources_metadata(existing_sources)
    return {"message": "Output successfully added to sources for recursive grounding!", "source": new_source}


# ==========================================
# GROUNDED CHAT STREAMING (LM STUDIO PROXY)
# ==========================================

@app.post("/api/chat")
async def chat_with_lm_studio(req: ChatRequest):
    base_url = (req.lm_studio_url or DEFAULT_LM_STUDIO_URL).rstrip("/")
    if not base_url.endswith("/v1"):
        base_url = f"{base_url}/v1"

    # Build grounded context from selected sources
    sources_meta = load_sources_metadata()
    grounded_chunks = []
    
    for sid in req.selected_source_ids:
        src = next((s for s in sources_meta if s["id"] == sid), None)
        if src:
            t_path = get_source_text_path(sid)
            if os.path.exists(t_path):
                with open(t_path, "r", encoding="utf-8") as f:
                    content = f.read()
                    grounded_chunks.append(f"=== SOURCE: {src['name']} (Type: {src.get('file_type', 'TXT')}) ===\n{content}\n=== END SOURCE: {src['name']} ===")

    context_str = "\n\n".join(grounded_chunks) if grounded_chunks else "No sources selected. Provide general assistance or ask the user to select sources."

    default_system = (
        "You are an expert analytical research assistant and document synthesizer (NotebookLM style).\n"
        "You have access to the user's grounded document sources below. Always reference and synthesize information "
        "from the provided sources accurately. Cite document names in brackets like [DocName.pdf] when answering factual questions.\n"
        "When requested to create reports, outlines, or presentations, structure your output with clean Markdown (headings, bullet points, concise executive summaries).\n\n"
        f"--- GROUNDED SOURCES CONTEXT ---\n{context_str}\n--- END CONTEXT ---"
    )
    
    system_content = req.system_prompt if req.system_prompt else default_system

    # Construct messages array
    messages = [{"role": "system", "content": system_content}]
    for m in req.messages:
        messages.append({"role": m.role, "content": m.content})

    # Generator for SSE / text streaming
    async def stream_generator():
        try:
            client = OpenAI(base_url=base_url, api_key=DEFAULT_API_KEY)
            
            # Fetch model if not specified
            model_to_use = req.model
            if not model_to_use:
                try:
                    models_list = client.models.list()
                    if models_list.data:
                        model_to_use = models_list.data[0].id
                    else:
                        model_to_use = "default"
                except Exception:
                    model_to_use = "local-model"

            response = client.chat.completions.create(
                model=model_to_use,
                messages=messages,
                temperature=req.temperature,
                stream=True
            )

            for chunk in response:
                if chunk.choices and len(chunk.choices) > 0:
                    delta = chunk.choices[0].delta
                    content_chunk = getattr(delta, "content", "") or ""
                    if content_chunk:
                        payload = json.dumps({"text": content_chunk})
                        yield f"data: {payload}\n\n"

            yield "data: [DONE]\n\n"

        except Exception as e:
            err_msg = (
                f"\n\n[Connection Error to LM Studio at {base_url}]: {str(e)}\n\n"
                "Please verify:\n"
                "1. LM Studio is running on your machine.\n"
                "2. The Local Server tab is ON at port 1234.\n"
                "3. A model is loaded into memory."
            )
            payload = json.dumps({"text": err_msg, "error": True})
            yield f"data: {payload}\n\n"
            yield "data: [DONE]\n\n"

    return StreamingResponse(stream_generator(), media_type="text/event-stream")


# ==========================================
# DOCUMENT EXPORT BUILDERS (DOCX & PPTX)
# ==========================================

def set_cell_background(cell, fill_hex):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), fill_hex)
    tcPr.append(shd)

def set_cell_margins(cell, top=100, bottom=100, left=150, right=150):
    tcPr = cell._tc.get_or_add_tcPr()
    tcMar = OxmlElement('w:tcMar')
    for m, val in [('top', top), ('bottom', bottom), ('left', left), ('right', right)]:
        node = OxmlElement(f'w:{m}')
        node.set(qn('w:w'), str(val))
        node.set(qn('w:type'), 'dxa')
        tcMar.append(node)
    tcPr.append(tcMar)

@app.post("/api/export/docx")
def export_docx(req: ExportDocxRequest):
    """
    Builds a styled Word (.docx) document from Markdown text using python-docx.
    Supports headings, lists, bold/italic, tables, callouts, and page formatting.
    """
    doc = Document()
    
    # Page setup (Standard Letter, 1-inch margins)
    section = doc.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    # Document Header / Title
    title_text = req.title or "Grounded Research Synthesis"
    title_p = doc.add_paragraph()
    title_run = title_p.add_run(title_text)
    title_run.font.name = "Arial"
    title_run.font.size = Pt(22)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(15, 23, 42) # Slate 900
    title_p.paragraph_format.space_after = Pt(4)

    # Subtitle / Metadata badge
    sub_p = doc.add_paragraph()
    sub_run = sub_p.add_run(f"Generated via Offline NotebookLM Studio • {datetime.now().strftime('%B %d, %Y')}")
    sub_run.font.name = "Arial"
    sub_run.font.size = Pt(9.5)
    sub_run.font.color.rgb = RGBColor(100, 116, 139) # Slate 500
    sub_p.paragraph_format.space_after = Pt(18)

    # Parse Markdown lines
    lines = req.markdown_content.split("\n")
    in_code_block = False
    code_lines = []

    for line in lines:
        stripped = line.strip()

        # Code block handling
        if stripped.startswith("```"):
            if in_code_block:
                # Flush code block
                table = doc.add_table(rows=1, cols=1)
                table.alignment = WD_TABLE_ALIGNMENT.CENTER
                cell = table.cell(0, 0)
                set_cell_background(cell, "F1F5F9")
                set_cell_margins(cell, top=140, bottom=140, left=180, right=180)
                cp = cell.paragraphs[0]
                c_run = cp.add_run("\n".join(code_lines))
                c_run.font.name = "Consolas"
                c_run.font.size = Pt(9)
                c_run.font.color.rgb = RGBColor(30, 41, 59)
                doc.add_paragraph().paragraph_format.space_after = Pt(6)
                code_lines = []
                in_code_block = False
            else:
                in_code_block = True
            continue

        if in_code_block:
            code_lines.append(line)
            continue

        if not stripped:
            continue

        # Heading 1 (# ...)
        if stripped.startswith("# "):
            h_text = stripped[2:].strip()
            p = doc.add_paragraph()
            run = p.add_run(h_text)
            run.font.name = "Arial"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(15, 23, 42)
            p.paragraph_format.space_before = Pt(16)
            p.paragraph_format.space_after = Pt(6)

        # Heading 2 (## ...)
        elif stripped.startswith("## "):
            h_text = stripped[3:].strip()
            p = doc.add_paragraph()
            run = p.add_run(h_text)
            run.font.name = "Arial"
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = RGBColor(30, 58, 138) # Indigo/Blue
            p.paragraph_format.space_before = Pt(12)
            p.paragraph_format.space_after = Pt(4)

        # Heading 3 (### ...)
        elif stripped.startswith("### "):
            h_text = stripped[4:].strip()
            p = doc.add_paragraph()
            run = p.add_run(h_text)
            run.font.name = "Arial"
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = RGBColor(51, 65, 85)
            p.paragraph_format.space_before = Pt(8)
            p.paragraph_format.space_after = Pt(3)

        # Blockquote (> ...)
        elif stripped.startswith(">"):
            quote_text = stripped[1:].strip()
            table = doc.add_table(rows=1, cols=1)
            cell = table.cell(0, 0)
            set_cell_background(cell, "F8FAFC")
            set_cell_margins(cell, top=100, bottom=100, left=160, right=160)
            qp = cell.paragraphs[0]
            q_run = qp.add_run(quote_text)
            q_run.font.name = "Georgia"
            q_run.font.size = Pt(10)
            q_run.font.italic = True
            q_run.font.color.rgb = RGBColor(51, 65, 85)
            doc.add_paragraph().paragraph_format.space_after = Pt(4)

        # Bullet lists (- ... or * ...)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            bullet_text = stripped[2:].strip()
            p = doc.add_paragraph(style='List Bullet')
            p.paragraph_format.space_before = Pt(2)
            p.paragraph_format.space_after = Pt(2)
            # Parse inline bold
            parts = bullet_text.split("**")
            for i, part in enumerate(parts):
                if not part:
                    continue
                r = p.add_run(part)
                r.font.name = "Arial"
                r.font.size = Pt(10)
                if i % 2 == 1:
                    r.font.bold = True
                    r.font.color.rgb = RGBColor(15, 23, 42)
                else:
                    r.font.color.rgb = RGBColor(51, 65, 85)

        # Standard Paragraph
        else:
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(2)
            p.paragraph_format.space_after = Pt(6)
            parts = stripped.split("**")
            for i, part in enumerate(parts):
                if not part:
                    continue
                r = p.add_run(part)
                r.font.name = "Arial"
                r.font.size = Pt(10)
                if i % 2 == 1:
                    r.font.bold = True
                    r.font.color.rgb = RGBColor(15, 23, 42)
                else:
                    r.font.color.rgb = RGBColor(51, 65, 85)

    # Save to memory buffer
    file_stream = io.BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)

    clean_filename = req.filename if req.filename and req.filename.endswith(".docx") else f"{req.filename or 'Synthesis'}.docx"
    return StreamingResponse(
        file_stream,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{clean_filename}"'}
    )


@app.post("/api/export/pptx")
def export_pptx(req: ExportPptxRequest):
    """
    Builds a structured presentation (.pptx) using python-pptx from markdown.
    Creates title slide + structured content slides with headers and bullet points.
    """
    prs = Presentation()
    # 16:9 widescreen layout
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme colors
    BG_COLOR = PptxRGBColor(248, 250, 252) # Slate 50
    CARD_COLOR = PptxRGBColor(255, 255, 255)
    ACCENT_COLOR = PptxRGBColor(37, 99, 235) # Blue 600
    TEXT_DARK = PptxRGBColor(15, 23, 42)     # Slate 900
    TEXT_MUTED = PptxRGBColor(100, 116, 139) # Slate 500

    # 1. Slide: Title Slide
    title_slide = prs.slides.add_slide(blank_layout)
    
    # Background shape
    bg_shape = title_slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = PptxRGBColor(15, 23, 42) # Deep Dark Slate for title
    bg_shape.line.fill.background()

    # Title box
    t_box = title_slide.shapes.add_textbox(PptxInches(1.2), PptxInches(2.2), PptxInches(11), PptxInches(3))
    tf = t_box.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "OFFLINE NOTEBOOKLM STUDIO"
    p0.font.name = "Arial"
    p0.font.size = PptxPt(13)
    p0.font.bold = True
    p0.font.color.rgb = PptxRGBColor(96, 165, 250) # Light blue badge
    p0.space_after = PptxPt(14)

    p1 = tf.add_paragraph()
    p1.text = req.title or "Executive Research & Synthesis"
    p1.font.name = "Arial"
    p1.font.size = PptxPt(36)
    p1.font.bold = True
    p1.font.color.rgb = PptxRGBColor(255, 255, 255)
    p1.space_after = PptxPt(16)

    p2 = tf.add_paragraph()
    p2.text = f"Local Document Grounding • LM Studio Integration • {datetime.now().strftime('%B %Y')}"
    p2.font.name = "Arial"
    p2.font.size = PptxPt(14)
    p2.font.color.rgb = PptxRGBColor(148, 163, 184)

    # 2. Parse markdown into slide sections
    # A new slide is triggered by ## or # or grouped every 4-5 bullet points
    raw_lines = req.markdown_content.split("\n")
    slides_data = []
    current_slide_title = ""
    current_bullets = []

    for line in raw_lines:
        s = line.strip()
        if not s:
            continue

        if s.startswith("# ") or s.startswith("## "):
            if current_slide_title or current_bullets:
                slides_data.append({
                    "title": current_slide_title or "Key Findings & Analysis",
                    "bullets": current_bullets if current_bullets else ["Synthesis overview"]
                })
            current_slide_title = s.lstrip("#").strip()
            current_bullets = []
        elif s.startswith("### "):
            # Subheading treated as strong bullet or title if none
            if not current_slide_title:
                current_slide_title = s.lstrip("#").strip()
            else:
                current_bullets.append(f"**{s.lstrip('#').strip()}**")
        elif s.startswith("- ") or s.startswith("* "):
            current_bullets.append(s[2:].strip())
        elif s.startswith(">"):
            current_bullets.append(f"Note: {s[1:].strip()}")
        else:
            # Regular paragraph
            if len(s) > 10:
                current_bullets.append(s)

    if current_slide_title or current_bullets:
        slides_data.append({
            "title": current_slide_title or "Summary & Action Items",
            "bullets": current_bullets if current_bullets else ["Analysis completed."]
        })

    # If no headings found, split bullets into slides of max 4 bullets
    if not slides_data:
        slides_data = [{"title": req.title or "Document Synthesis", "bullets": ["No structured bullet points provided."]}]

    # Create Content Slides
    for idx, s_info in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        
        # Background
        bg = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

        # Content Card shape
        card = slide.shapes.add_shape(
            pptx.enum.shapes.MSO_SHAPE.ROUNDED_RECTANGLE,
            PptxInches(0.9), PptxInches(0.8), PptxInches(11.533), PptxInches(5.9)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = PptxRGBColor(226, 232, 240)
        card.line.width = PptxPt(1)

        # Slide Title
        tbox = slide.shapes.add_textbox(PptxInches(1.4), PptxInches(1.2), PptxInches(10.5), PptxInches(1.0))
        tf = tbox.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = s_info["title"]
        p_t.font.name = "Arial"
        p_t.font.size = PptxPt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK

        # Top Accent Line
        accent_line = slide.shapes.add_shape(
            pptx.enum.shapes.MSO_SHAPE.RECTANGLE,
            PptxInches(1.4), PptxInches(2.2), PptxInches(2.5), PptxInches(0.04)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = ACCENT_COLOR
        accent_line.line.fill.background()

        # Slide Body Bullets
        bbox = slide.shapes.add_textbox(PptxInches(1.4), PptxInches(2.4), PptxInches(10.5), PptxInches(3.8))
        btf = bbox.text_frame
        btf.word_wrap = True

        bullets = s_info["bullets"][:6] # Limit to 6 items per slide
        for b_idx, bullet in enumerate(bullets):
            p_b = btf.paragraphs[0] if b_idx == 0 else btf.add_paragraph()
            clean_b = bullet.replace("**", "").replace("`", "")
            p_b.text = f"•   {clean_b}"
            p_b.font.name = "Arial"
            p_b.font.size = PptxPt(14)
            p_b.font.color.rgb = PptxRGBColor(51, 65, 85)
            p_b.space_after = PptxPt(12)

        # Slide number footer
        fbox = slide.shapes.add_textbox(PptxInches(10.5), PptxInches(6.2), PptxInches(1.8), PptxInches(0.4))
        ftf = fbox.text_frame
        fp = ftf.paragraphs[0]
        fp.alignment = PP_ALIGN.RIGHT
        fp.text = f"Slide {idx + 2} of {len(slides_data) + 1}"
        fp.font.name = "Arial"
        fp.font.size = PptxPt(10)
        fp.font.color.rgb = TEXT_MUTED

    # Save to memory buffer
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)

    clean_filename = req.filename if req.filename and req.filename.endswith(".pptx") else f"{req.filename or 'Presentation'}.pptx"
    return StreamingResponse(
        file_stream,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{clean_filename}"'}
    )


# ==========================================
# STATIC FRONTEND SERVING
# ==========================================

# Mount static assets if dist exists
dist_dir = os.path.join(BASE_DIR, "dist")
dist_assets = os.path.join(dist_dir, "assets")
if os.path.exists(dist_assets):
    app.mount("/assets", StaticFiles(directory=dist_assets), name="assets")

# Serve frontend on both / and /app
@app.get("/", response_class=HTMLResponse)
@app.get("/app", response_class=HTMLResponse)
def serve_app_html():
    dist_index_path = os.path.join(BASE_DIR, "dist", "index.html")
    standalone_path = os.path.join(BASE_DIR, "standalone_index.html")
    
    if os.path.exists(dist_index_path):
        with open(dist_index_path, "r", encoding="utf-8") as f:
            return f.read()
    elif os.path.exists(standalone_path):
        with open(standalone_path, "r", encoding="utf-8") as f:
            return f.read()
    
    index_path = os.path.join(BASE_DIR, "index.html")
    if os.path.exists(index_path):
        with open(index_path, "r", encoding="utf-8") as f:
            return f.read()
    return "<h1>Offline NotebookLM Studio</h1><p>FastAPI Server is running.</p>"

@app.get("/api/health")
def health_check():
    return {
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "sources_count": len(load_sources_metadata()),
        "lm_studio_endpoint": DEFAULT_LM_STUDIO_URL
    }

if __name__ == "__main__":
    print("=" * 60)
    print("🚀 Starting Offline NotebookLM Studio on http://localhost:8000")
    print(f"📡 Configured for LM Studio at {DEFAULT_LM_STUDIO_URL}")
    print("=" * 60)
    uvicorn.run(app, host="0.0.0.0", port=8000)
